import tkinter as tk
from tkinter import ttk, scrolledtext, messagebox, filedialog, simpledialog
import os
import platform
import subprocess
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN, MSO_AUTO_SIZE
import sys
import copy
import requests
from datetime import date
import re
from tkinter import simpledialog, messagebox
from tkinter import font as tkfont
# --- Constantes de Configuração ---
LITURGIA_API_URL = "https://api-liturgia-diaria.vercel.app/"
LARGURA_SLIDE = Inches(16)
ALTURA_SLIDE = Inches(9)
MARGEM_TEXTO = Inches(0.01)
DEFAULT_TAMANHO_FONTE_MUSICA_REFRAO = Pt(80)
DEFAULT_TAMANHO_FONTE_MUSICA_VERSO = Pt(80)
DEFAULT_TAMANHO_FONTE_ACLAMACAO = Pt(70)
DEFAULT_TAMANHO_FONTE_ANTIFONA = Pt(66)
DEFAULT_TAMANHO_FONTE_LEITURA_TITULO_AMARELO = Pt(90)
DEFAULT_TAMANHO_FONTE_LEITURA_TEXTO_BRANCO = Pt(90)
DEFAULT_TAMANHO_FONTE_PALAVRA = Pt(80)
TAMANHO_TITULO_PARTE = Pt(96)
TAMANHO_FONTE_TITULO_INICIAL = Pt(90)
LINHAS_POR_SLIDE_VERSO = 3
LINHAS_POR_SLIDE_LEITURA = 5
LINHAS_POR_SLIDE_PALAVRA = 6
NOME_FONTE_PADRAO = 'Arial'
# Conversões aproximadas para medição em pixels (para paginador inteligente)
DPI_SCREEN = 96
SLIDE_PX_WIDTH = int(16 * DPI_SCREEN)
SLIDE_PX_HEIGHT = int(9 * DPI_SCREEN)
CAIXA_MARGIN_PX = 24
# Paleta/estilos de UI (somente aparência; sem alterar funcionalidades)
UI_COLOR_PRIMARY = '#FFC000'   # Amarelo (similar aos títulos)
UI_COLOR_BG = '#F7F7F7'
UI_COLOR_ACCENT = '#1E90FF'
UI_FONT_BASE = ('Arial', 10)
UI_FONT_BOLD = ('Arial', 10, 'bold')
FONTES_COMUNS_PPT = sorted([
    "Arial", "Arial Black", "Arial Narrow", "Bahnschrift", "Calibri", "Calibri Light", "Cambria", "Cambria Math", "Candara", "Candara Light", "Century", "Century Gothic", "Century Schoolbook", "Comic Sans MS", "Consolas", "Constantia", "Corbel", "Corbel Light", "Courier New", "Ebrima", "Franklin Gothic Medium", "Franklin Gothic Book", "Gabriola", "Gadugi", "Georgia", "Gill Sans MT", "Impact", "Ink Free", "Leelawadee UI", "Lucida Console", "Lucida Sans Unicode", "Malgun Gothic", "Marlett", "Microsoft Himalaya", "Microsoft JhengHei", "Microsoft JhengHei UI", "Microsoft New Tai Lue", "Microsoft PhagsPa", "Microsoft Sans Serif", "Microsoft Tai Le", "Microsoft YaHei", "Microsoft YaHei UI", "Microsoft Yi Baiti", "MingLiU-ExtB", "PMingLiU-ExtB", "MingLiU_HKSCS-ExtB", "Mongolian Baiti", "Montserrat", "MS Gothic", "MS UI Gothic", "MS PGothic", "MV Boli", "Myanmar Text", "Nirmala UI", "Palatino Linotype", "Rockwell", "Segoe Print", "Segoe Script", "Segoe UI", "Segoe UI Black", "Segoe UI Emoji", "Segoe UI Historic", "Segoe UI Semibold", "Segoe UI Semilight", "Segoe UI Symbol", "SimSun", "NSimSun", "SimSun-ExtB", "Sitka Banner", "Sitka Display", "Sitka Heading", "Sitka Small", "Sitka Subheading", "Sitka Text", "Sylfaen", "Symbol", "Tahoma", "Times New Roman", "Trebuchet MS", "Verdana", "Webdings", "Wingdings", "Wingdings 2", "Wingdings 3"
])
COR_REFRAO = RGBColor(255, 192, 0); COR_VERSO = RGBColor(255, 255, 255)
COR_TITULO = RGBColor(255, 192, 0); COR_FUNDO_PRETO = RGBColor(0, 0, 0)

DEFAULT_TEXTS_ORIGINAL = {
     "Entrada": {"titulo": "CANTO DE ENTRADA", "refrao": [], "versos": []},
     "Ato Penitencial": {"titulo": "ATO PENITENCIAL", "refrao": [], "versos": []},
     "Glória": {"titulo": "GLÓRIA", "refrao": [], "versos": []},
     "Palavra": {"titulo": "PALAVRA", "texto": ["DESÇA COMO A CHUVA A TUA","PALAVRA. QUE SE ESPALHE","COMO ORVALHO. COMO","CHUVISCO NA RELVA. COMO","AGUACEIRO NA GRAMA.","AMÉM!"]},
     "1ª Leitura": {"titulo_amarelo": ["PRIMEIRA LEITURA"], "texto_branco": ["Isaias 55,10-11"] },
     "Salmo": {"titulo_amarelo": ["SALMO 64 (65)"], "texto_branco": ["- A semente caiu em terra boa e deu fruto."] },
     "2ª Leitura": {"titulo_amarelo": ["SEGUNDA LEITURA"], "texto_branco": ["Romanos 8,18-23"] },
     "Aclamação": {"titulo": "ACLAMAÇÃO DO EVANGELHO", "aclamacao_texto": ["Aleluia, Aleluia, Aleluia!"], "antifona_texto": ["Tua Palavra é a luz do caminho", "A lâmpada para os meus pés, Senhor!", "(Mt 13,1-23)"]},
     "Oferendas": {"titulo": "PREPARAÇÃO DAS OFERENDAS", "refrao": [], "versos": []},
     "Comunhão": {"titulo": "COMUNHÃO", "refrao": [], "versos": []},
}

# Conteúdos de Orações Eucarísticas (cada item é uma "oração" com suas linhas; linhas vazias indicam quebra manual de slide)
ORACOES_EUCARISTICAS = {
    "ORAÇÃO EUCARÍSTICA Nº 03 — DIVERSAS CIRCUNSTÂNCIAS": [
        "SANTO",
        "",
        "Bendito o vosso Filho, presente entre nós!",
        "",
        "Enviai o vosso Espírito Santo!",
        "",
        "Anunciamos, Senhor, a vossa morte e proclamamos a vossa ressurreição. Vinde, Senhor Jesus!",
        "",
        "Aceitai, ó Senhor, a nossa oferta!",
        "",
        "O Espírito nos una num só corpo!",
        "",
        "Confirmai na unidade a vossa Igreja",
        "",
        "Ajudai-nos a criar um mundo novo!",
        "",
        "Ajudai-nos a criar um mundo novo!",
    ],
    "ORAÇÃO EUCARÍSTICA Nº 02": [
        "SANTO",
        "",
        "Enviai o vosso espirito santo",
        "",
        "Anunciamos, Senhor, a vossa morte e proclamamos a vossa ressurreição. ",
        "Vinde, Senhor Jesus!",
        "",
        "Aceitai, ó Senhor, a nossa oferta!",
        "",
        "O Espirito nos una num só corpo",
        "",
        "Lembrai-Vos, ó Pai, da Vossa Igreja! ",
        "",
        "Concedei-lhes, ó Senhor, a luz eterna!",
    ],
    "ORAÇÃO EUCARÍSTICA Nº 01": [
        "SANTO",
        "",
        "Abençoai nossa oferenda, ó Senhor!",
        "",
        "Lembrai-vos, ó Pai dos vossos filhos!",
        "",
        "Em comunhão com vossos Santos, vos louvamos!",
        "",
        "Enviai o vosso Espírito Santo!",
        "",
        "Anunciamos, Senhor, a vossa morte e proclamamos a vossa ressurreição. Vinde, Senhor Jesus!",
        "",
        "Salvador do mundo, salvai-nos, vós que nos libertastes pela cruz e ressureição!",
        "",
        "Aceitai, ó Senhor, a nossa oferta!",
        "",
        "O Espírito nos una num só corpo!",
        "",
        "Concedei-lhes, ó Senhor, a luz eterna!",
    ],
    "ORAÇÃO EUCARÍSTICA Nº 05": [
        "SANTO",
        "",
        "Mandai vosso Espírito Santo!",
        "",
        "Toda vez que se come deste pão, toda vez que se bebemos deste vinho, recordamos a paixão de Jesus Cristo e  ficamos esperando sua vinda!",
        "",
        "Recebei, ó Senhor, a nossa oferta!",
        "",
        "O Espírito nos una num só corpo.",
        "",
        "Caminhamos na estrada de Jesus.",
        "",
        "Lembrai- vos ,ó Pai , da vossa Igreja!",
        "",
        "Esperamos entrar na vida eterna.",
        "",
        "A todos daí a luz que não se apaga.",
    ],
    "ORAÇÃO EUCARÍSTICA Nº 04": [
        "Alegrai-nos, ó Pai, com Vossa luz!",
        "",
        "SANTO",
        "",
        "Socorrei, com bondade, os que Vos buscam!",
        "",
        "Por amor nos enviastes Vosso Filho!",
        "",
        "Jesus Cristo deu-nos a vida por Sua morte!",
        "",
        "Santificai-nos pelo dom do vosso Espírito!",
        "",
        "Santificai Nossa oferenda pelo Espírito!",
        "",
        "Salvador do mundo, salvai-nos, Vós que nos libertastes pela cruz e ressurreição!",
        "",
        "Recebei, ó Senhor, a nossa oferta!",
        "",
        "Fazei de nos um sacrifício de louvor!",
        "",
        "Lembrai-Vos, ó Pai, dos Vossos filhos!",
        "",
        "A todos saciai com Vossa glória!",
        "",
        "Concedei-nos o convívio dos eleitos!",
    ],
    "ORAÇÃO EUCARÍSTICA Nº 03 (VARIAÇÃO)": [
        "SANTO",
        "",
        "Enviai o vosso Espírito Santo!",
        "",
        "Anunciamos, Senhor, a vossa morte e proclamamos a vossa ressurreição. Vinde, Senhor Jesus!",
        "",
        "Aceitai, ó Senhor, a nossa oferta!",
        "",
        "O Espírito nos una num só corpo!",
        "",
        "Fazei de nós uma perfeita oferenda!",
        "",
        "Lembrai-Vos, ó Pai, da Vossa Igreja!",
    ],
    "ORAÇÃO EUCARÍSTICA VI-C": [
        "O vosso filho permaneça entre nós!",
        "",
        "Mandai o vosso Espírito Santo!",
        "",
        "Anunciamos, Senhor a vossa morte e proclamamos a vossa ressureição. Vinde, Senhor Jesus!",
        "",
        "Aceitai, ó Senhor, a nossa oferta!",
        "",
        "O vosso espirito nos uma num só corpo!",
        "",
        "Caminhamos no amor e na alegria!",
        "",
        "Concedei-lhes, ó Senhor, a luz eterna!",
    ],
    "ORAÇÃO EUCARÍSTICA 6B": [
        "O vosso filho permaneça entre nós!",
        "",
        "Mandai o vosso Espírito Santo!",
        "",
        "Salvador do mundo, salvai-nos, vós que nos libertastes pela cruz e ressurreição.",
        "",
        "Aceitai, ó Senhor,",
        "a nossa oferta!",
        "",
        "Tornai viva nossa fé, nossa esperança!",
        "",
        "Concedei-lhes,",
        "ó Senhor,",
        "a luz eterna! ",
    ],
    "ORAÇÃO EUCARÍSTICA 6A": [
        "O Vosso filho permaneça entre nós!",
        "",
        "Mandai o vosso Espírito Santo!",
        "",
        "Todas as vezes que comemos deste pão e bebemos deste cálice, anunciamos Senhor, a vossa morte, enquanto esperamos a vossa vinda!",
        "",
        "Aceitai, ó Senhor, a nossa oferta!",
        "",
        "Confirmai na caridade o vosso povo!",
        "",
        "Concedei-lhes ó Senhor, a luz eterna !",
    ],
    "ORAÇÃO EUCARÍSTICA 7": [
        " Como é grande, ó Pai, a vossa misericórdia!",
        "",
        " Como é grande, ó Pai, a vossa misericórdia!",
        "",
        " Como é grande, ó Pai, a vossa misericórdia!",
        "",
        "Anunciamos, Senhor, a vossa morte e proclamamos a vossa ressurreição. Vinde, Senhor Jesus!",
        "",
        "Esperamos, ó Cristo, vossa vinda gloriosa!",
        "",
        "Esperamos, ó Cristo, vossa vinda gloriosa!",
        "",
        "Esperamos, ó Cristo, vossa vinda gloriosa!",
    ],
    "ORAÇÃO EUCARÍSTICA VI-D (6-D)": [
        "BENDITO O VOSSO FILHO, PRESENTE ENTRE NÓS!",
        "",
        "ENVIAI O VOSSO ESPÍRITO SANTO!",
        "",
        "ANUNCIAMOS, SENHOR A VOSSA MORTE E PROCLAMAMOS A VOSSA RESSUREIÇÃO. VINDE, SENHOR JESUS!",
        "",
        "ACEITAI, Ó SENHOR, A NOSSA OFERTA!",
        "",
        "O ESPÍRITO NOS UNA NUM SÓ CORPO!",
        "",
        "CONFIRMAI NA UNIDADE A VOSSA IGREJA!",
        "",
        "AJUDAI-NOS A CRIAR UM MUNDO NOVO!",
        "",
        "CONCEDEI-LHES, Ó SENHOR, A LUZ ETERNA!",
    ],
    "ORAÇÃO EUCARÍSTICA Nº 03 — DIVERSAS CIRCUNSTÂNCIAS (VARIAÇÃO 2)": [
        "SANTO",
        "",
        "Bendito o vosso Filho, presente entre nós!",
        "",
        "Enviai o vosso Espírito Santo!",
        "",
        "Anunciamos, Senhor, a vossa morte e proclamamos a vossa ressurreição. Vinde, Senhor Jesus!",
        "",
        "Aceitai, ó Senhor, a nossa oferta!",
        "",
        "O Espírito nos una num só corpo!",
        "",
        "Confirmai na unidade a vossa Igreja",
        "",
        "Ajudai-nos a criar um mundo novo!",
        "",
        "Concedei-lhes, ó senhor, a luz eterna!",
    ],
    "ORAÇÃO EUCARÍSTICA Nº 08 — SOBRE RECONCILIAÇÃO II": [
        "FAZEI-NOS,Ó PAI, INSTRUMENTOS DE VOSSA PAZ!",
        "",
        "SANTO",
        "",
        "FAZEI-NOS,Ó PAI, INSTRUMENTOS DE VOSSA PAZ!",
        "",
        "ANUNCIAMOS, SENHOR, A VOSSA MORTE, E PROCLAMAMOS A VOSSA RESSURREIÇÃO. VINDE, SENHOR JESUS!",
        "",
        "GLÓRIA E LOUVOR AO PAI, QUE EM CRISTO NOS RECONCILIOU!",
        "",
        "GLÓRIA E LOUVOR AO PAI, QUE EM CRISTO NOS RECONCILIOU!",
        "",
        "GLÓRIA E LOUVOR AO PAI, QUE EM CRISTO NOS RECONCILIOU!",
        "",
        "GLÓRIA E LOUVOR AO PAI, QUE EM CRISTO NOS RECONCILIOU!",
    ],
}
TEXTO_CREDO = [ "CREIO EM DEUS PAI TODO PODEROSO,", "CRIADOR DO CÉU E DA TERRA.", "E EM JESUS CRISTO, SEU ÚNICO FILHO,", "NOSSO SENHOR,", "QUE FOI CONCEBIDO PELO PODER DO ESPÍRITO SANTO;", "NASCEU DA VIRGEM MARIA;", "PADECEU SOB PÔNCIO PILATOS,", "FOI CRUCIFICADO, MORTO E SEPULTADO.", "DESCEU À MANSÃO DOS MORTOS;", "RESSUSCITOU AO TERCEIRO DIA;", "SUBIU AOS CÉUS, ESTÁ SENTADO À DIREITA", "DE DEUS PAI TODO-PODEROSO,", "DONDE HÁ DE VIR A JULGAR OS VIVOS E OS MORTOS.", "CREIO NO ESPÍRITO SANTO,", "NA SANTA IGREJA CATÓLICA,", "NA COMUNHÃO DOS SANTOS,", "NA REMISSÃO DOS PECADOS,", "NA RESSURREIÇÃO DA CARNE,", "NA VIDA ETERNA.", "AMÉM." ]
TEXTO_ORACAO_SANTA_LUZIA = [ "Ó VIRGEM ADMIRÁVEL.", "CHEIA DE FIRMEZA E DE", "CONSTÂNCIA, QUE NEM", "AS POMPAS HUMANAS", "PUDERAM SEDUZIR,", "NEM AS PROMESSAS,", "NEM AS AMEAÇAS,", "NEM A FORÇA BRUTA", "PUDERAM ABALAR,", "PORQUE SOUBESTES SER", "O TEMPLO VIVO DO", "DIVINO ESPÍRITO SANTO.", "O MUNDO CRISTÃO VOS", "PROCLAMOU ADVOGADA", "DA LUZ DOS NOSSOS", "OLHOS. DEFENDEI-NOS,", "POIS, DE TODA MOLÉSTIA", "QUE POSSA PREJUDICAR", "A NOSSA VISTA.", "ALCANÇAI-NOS A LUZ", "SOBRENATURAL DA FÉ,", "ESPERANÇA E CARIDADE", "PARA QUE NOS", "DESAPEGUEMOS", "DAS COISAS MATERIAIS", "E TERRESTRES", "E TENHAMOS A FORÇA", "PARA VENCER O INIMIGO", "E ASSIM POSSAMOS", "CONTEMPLAR-VOS NA", "GLÓRIA CELESTE. AMÉM." ]

def adiciona_texto_com_divisao(prs, layout, linhas_originais, cor, tamanho_fonte, font_name, bold_state, italic_state, max_linhas, use_auto_size=True):
    if not linhas_originais or all(not s or s.isspace() for s in linhas_originais): return False
    linhas_validas = [linha for linha in linhas_originais if linha and not linha.isspace()]
    if not linhas_validas: return False
    linhas_restantes = linhas_validas[:]; slides_criados = 0
    while linhas_restantes:
        linhas_para_este_bloco = linhas_restantes[:max_linhas]; linhas_restantes = linhas_restantes[max_linhas:]
        texto_bloco_continuo = " ".join(linhas_para_este_bloco)
        if not texto_bloco_continuo.strip(): continue
        slide = prs.slides.add_slide(layout); slides_criados += 1
        esquerda = MARGEM_TEXTO; topo = MARGEM_TEXTO; largura = LARGURA_SLIDE - (2 * MARGEM_TEXTO); altura = ALTURA_SLIDE - (2 * MARGEM_TEXTO)
        caixa_texto = slide.shapes.add_textbox(esquerda, topo, largura, altura)
        frame_texto = caixa_texto.text_frame; frame_texto.clear(); frame_texto.word_wrap = True; frame_texto.vertical_anchor = MSO_ANCHOR.MIDDLE
        if use_auto_size: frame_texto.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        else: frame_texto.auto_size = MSO_AUTO_SIZE.NONE
        p = frame_texto.add_paragraph(); p.text = texto_bloco_continuo; p.alignment = PP_ALIGN.CENTER;
        p.font.name = font_name; p.font.size = tamanho_fonte; p.font.color.rgb = cor; p.font.bold = bold_state; p.font.italic = italic_state;
        try:
            caixa_texto.left = esquerda; caixa_texto.top = topo; caixa_texto.width = largura; caixa_texto.height = altura
            frame_texto.margin_bottom = Inches(0.05); frame_texto.margin_top = Inches(0.05); frame_texto.margin_left = Inches(0.1); frame_texto.margin_right = Inches(0.1)
        except Exception as e_resize: print(f"Aviso: resize caixa: {e_resize}")
    return slides_criados > 0

# ====== Paginador Inteligente (quebra por largura/altura, sem alterar funcionalidades) ======
def _hyphenate_long_word(word: str, font_obj, max_width_px: int) -> list:
    if font_obj.measure(word) <= max_width_px:
        return [word]
    # quebra simples por hífen quando a palavra excede a largura
    parts = []
    current = ''
    for ch in word:
        tentative = current + ch
        if font_obj.measure(tentative + '-') > max_width_px and current:
            parts.append(current + '-')
            current = ch
        else:
            current = tentative
    if current:
        parts.append(current)
    return parts

def _wrap_line_to_width(line: str, font_obj, max_width_px: int) -> list:
    if not line:
        return ['']
    words = line.split()
    wrapped = []
    current = ''
    for word in words:
        candidate = (current + ' ' + word).strip()
        if font_obj.measure(candidate) <= max_width_px:
            current = candidate
            continue
        # palavra sozinha estoura; tenta hifenizar
        if not current:
            for part in _hyphenate_long_word(word, font_obj, max_width_px):
                if not current:
                    current = part
                else:
                    if font_obj.measure((current + ' ' + part).strip()) <= max_width_px:
                        current = (current + ' ' + part).strip()
                    else:
                        wrapped.append(current)
                        current = part
            continue
        # fecha linha e começa nova com a palavra
        wrapped.append(current)
        current = word
    if current or not wrapped:
        wrapped.append(current)
    return wrapped

def _paginate_lines_by_area(lines: list, font_name: str, pt_size: int, area_width_px: int, area_height_px: int, line_spacing: float = 1.15) -> list:
    font_obj = tkfont.Font(family=font_name, size=pt_size, weight='bold')
    line_px = int(font_obj.metrics('linespace') * line_spacing)
    max_lines_per_page = max(1, area_height_px // max(1, line_px))
    wrapped_lines = []
    for original in lines:
        wrapped_lines.extend(_wrap_line_to_width(original, font_obj, area_width_px))
    pages = []
    while wrapped_lines:
        page = wrapped_lines[:max_lines_per_page]
        pages.append(page)
        wrapped_lines = wrapped_lines[max_lines_per_page:]
    return pages

def adiciona_texto_com_paginacao_inteligente(prs, layout, linhas_originais, cor, tamanho_fonte_pt, font_name, bold_state, italic_state, area_width_px=None, area_height_px=None):
    if not linhas_originais:
        return False
    if area_width_px is None:
        area_width_px = max(300, SLIDE_PX_WIDTH - 2 * CAIXA_MARGIN_PX)
    if area_height_px is None:
        area_height_px = max(200, SLIDE_PX_HEIGHT - 2 * CAIXA_MARGIN_PX)
    paginas = _paginate_lines_by_area(linhas_originais, font_name, int(tamanho_fonte_pt.pt), area_width_px, area_height_px)
    if not paginas:
        return False
    conteudo = False
    for page_lines in paginas:
        slide = prs.slides.add_slide(layout)
        esquerda = MARGEM_TEXTO; topo = MARGEM_TEXTO; largura = LARGURA_SLIDE - (2 * MARGEM_TEXTO); altura = ALTURA_SLIDE - (2 * MARGEM_TEXTO)
        caixa_texto = slide.shapes.add_textbox(esquerda, topo, largura, altura)
        frame_texto = caixa_texto.text_frame
        frame_texto.clear(); frame_texto.word_wrap = True; frame_texto.vertical_anchor = MSO_ANCHOR.MIDDLE
        frame_texto.auto_size = MSO_AUTO_SIZE.NONE
        for idx, ln in enumerate(page_lines):
            if idx == 0:
                p = frame_texto.add_paragraph()
            else:
                p = frame_texto.add_paragraph()
            p.text = ln
            p.alignment = PP_ALIGN.CENTER
            p.font.name = font_name; p.font.size = tamanho_fonte_pt; p.font.color.rgb = cor; p.font.bold = bold_state; p.font.italic = italic_state
        try:
            caixa_texto.left = esquerda; caixa_texto.top = topo; caixa_texto.width = largura; caixa_texto.height = altura
            frame_texto.margin_bottom = Inches(0.05); frame_texto.margin_top = Inches(0.05); frame_texto.margin_left = Inches(0.1); frame_texto.margin_right = Inches(0.1)
        except Exception:
            pass
        conteudo = True
    return conteudo

class MassSlideGeneratorApp:
    def __init__(self, master):
        self.master = master
        
        self.style = ttk.Style()
        available_themes = self.style.theme_names()
        if 'clam' in available_themes: self.style.theme_use('clam')
        elif 'alt' in available_themes: self.style.theme_use('alt')
        
        self._aplicar_estilos_globais()

        master.title("Slides To My Church v35.2 - UI Clássica / Métodos Completos") 
        master.geometry("1150x950") 
        try:
            # Ícone da aplicação (não altera comportamento, apenas visual)
            icon_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'icone.ico')
            if os.path.exists(icon_path):
                master.iconbitmap(icon_path)
        except Exception:
            pass
        try:
            master.minsize(900, 700)
            self._centralizar_janela(master)
        except Exception:
            pass

        self.DEFAULT_TEXTS = copy.deepcopy(DEFAULT_TEXTS_ORIGINAL) 

        manage_sections_outer_frame = ttk.Frame(master, padding=(12,12,12,6))
        manage_sections_outer_frame.pack(fill="x")
        manage_sections_frame = ttk.LabelFrame(manage_sections_outer_frame, text="Gerenciar Seções", padding="10")
        manage_sections_frame.pack(fill="x")
        
        btn_add = ttk.Button(manage_sections_frame, text="Adicionar Nova Seção ", command=self.dialogo_adicionar_secao)
        btn_add.pack(side="left", padx=(0,8), pady=2)
        self._adicionar_tooltip(btn_add, "Cria uma nova seção musical personalizada")
        self.remove_button = ttk.Button(manage_sections_frame, text="Remover Seção", command=self.remover_secao_selecionada, state="disabled")
        self.remove_button.pack(side="left", padx=5, pady=2)
        self._adicionar_tooltip(self.remove_button, "Remove a aba atualmente selecionada")
        self.move_up_button = ttk.Button(manage_sections_frame, text="Mover Para Cima ↑", command=lambda: self.mover_secao_selecionada(-1), state="disabled")
        self.move_up_button.pack(side="left", padx=5, pady=2)
        self._adicionar_tooltip(self.move_up_button, "Move a aba selecionada uma posição acima")
        self.move_down_button = ttk.Button(manage_sections_frame, text="Mover Para Baixo ↓", command=lambda: self.mover_secao_selecionada(1), state="disabled")
        self.move_down_button.pack(side="left", padx=5, pady=2)
        self._adicionar_tooltip(self.move_down_button, "Move a aba selecionada uma posição abaixo")
                
        self.fetch_liturgia_button = ttk.Button(manage_sections_frame, text="Buscar Liturgia por Data", command=self.fetch_liturgia_por_data)
        self.fetch_liturgia_button.pack(side="left", padx=5, pady=2)
        self._adicionar_tooltip(self.fetch_liturgia_button, "Importa leituras da API por uma data específica")
        self.config_oracoes_button = ttk.Button(manage_sections_frame, text="Orações Eucarísticas", command=self.configurar_oracoes_eucaristicas_dialog)
        self.config_oracoes_button.pack(side="left", padx=5, pady=2)
        self._adicionar_tooltip(self.config_oracoes_button, "Seleciona e estiliza as Orações Eucarísticas")

        ttk.Separator(master, orient='horizontal').pack(fill='x', padx=12, pady=(4, 6))

        title_outer_frame = ttk.Frame(master, padding=(12,6,12,6))
        title_outer_frame.pack(fill="x")
        title_frame = ttk.LabelFrame(title_outer_frame, text="Título Inicial da Apresentação", padding="10")
        title_frame.pack(fill="x")
        self.initial_title_widget = scrolledtext.ScrolledText(title_frame, height=3, width=90, wrap=tk.WORD, font=UI_FONT_BASE); 
        self.initial_title_widget.pack(fill="x", expand=True, pady=(2, 0))
        self.initial_title_widget.insert(tk.END, "DOMINGO DA\nQUARESMA")

        ttk.Separator(master, orient='horizontal').pack(fill='x', padx=12, pady=(0, 6))
        notebook_frame = ttk.Frame(master, padding=(12,6,12,6)) 
        notebook_frame.pack(expand=True, fill="both")
        self.notebook = ttk.Notebook(notebook_frame)
        self.notebook.bind("<<NotebookTabChanged>>", self._atualizar_estado_botoes_controle_abas) 
        self.widgets_gui = {}
        
        self.ordem_gui_inicial = [ "Entrada", "Ato Penitencial", "Glória", "Palavra", "1ª Leitura", "Salmo", "2ª Leitura", "Aclamação", "Oferendas", "Comunhão" ]
        
        self.ordem_geracao_dinamica = [] 

        for nome_parte in self.ordem_gui_inicial:
            if nome_parte in self.DEFAULT_TEXTS:
                 self._criar_aba_secao(nome_parte, tipo_override=None, inserir_em_posicao=-1, reconstruir_ordem=False) 
        
        self._reconstruir_ordem_geracao_dinamica() 
        self._atualizar_estado_botoes_controle_abas() 

        self.notebook.pack(expand=True, fill="both") 
        
        ttk.Separator(master, orient='horizontal').pack(fill='x', padx=12, pady=(0, 6))
        bottom_controls_frame = ttk.Frame(master, padding=(12,6,12,12))
        bottom_controls_frame.pack(fill="x", side="bottom")
        self.status_label = ttk.Label(bottom_controls_frame, text="Pronto."); 
        self.status_label.pack(side="left", padx=(0,10), pady=5)
        self.generate_button = ttk.Button(bottom_controls_frame, text="Gerar PowerPoint", command=self.gerar_apresentacao, style="Accent.TButton"); 
        self.generate_button.pack(side="right", pady=5)
        self._adicionar_tooltip(self.generate_button, "Gera o arquivo .pptx com os slides da celebração")

    def _aplicar_estilos_globais(self):
        try:
            self.style.configure("TFrame", background=UI_COLOR_BG)
        except Exception:
            pass
        self.style.configure("TLabel", padding=2, font=UI_FONT_BASE)
        self.style.configure("TLabelframe", padding=8)
        self.style.configure("TLabelframe.Label", font=UI_FONT_BOLD, foreground='black')
        self.style.configure("TButton", padding=(8,4))
        self.style.configure("TCombobox", padding=2)
        self.style.configure("TSpinbox", padding=2)
        self.style.configure("TCheckbutton", padding=2)
        try:
            self.style.configure("Accent.TButton", font=UI_FONT_BOLD)
        except tk.TclError:
            pass

    def _centralizar_janela(self, window):
        window.update_idletasks()
        w = window.winfo_width(); h = window.winfo_height()
        sw = window.winfo_screenwidth(); sh = window.winfo_screenheight()
        x = int((sw - w) / 2); y = int((sh - h) / 2)
        window.geometry(f"{w}x{h}+{x}+{y}")

    def _adicionar_tooltip(self, widget, text):
        # Tooltip simples, não intrusivo
        tip = tk.Toplevel(widget)
        tip.withdraw(); tip.overrideredirect(True)
        tip.wm_attributes("-topmost", True)
        lbl = ttk.Label(tip, text=text, relief='solid', borderwidth=1, padding=(6,3))
        lbl.pack()
        def enter(_):
            try:
                x = widget.winfo_rootx() + 20
                y = widget.winfo_rooty() + 20
                tip.geometry(f"+{x}+{y}"); tip.deiconify()
            except Exception:
                pass
        def leave(_):
            tip.withdraw()
        widget.bind("<Enter>", enter)
        widget.bind("<Leave>", leave)


        # Busca inicial silenciosa da liturgia do dia
        try:
            self.master.after(300, self.carregar_liturgia_hoje_silencioso)
        except Exception:
            pass

        # Estado das Orações Eucarísticas
        self.eucaristica_selecionar_todas = tk.BooleanVar(value=True)
        self.eucaristica_escolhidas = set()
        self.eucaristica_font_size_pt = 90


    def _abrir_dialogo_data(self):
        dialog = tk.Toplevel(self.master)
        dialog.title("Escolher data")
        dialog.transient(self.master)
        dialog.grab_set()
        dialog.resizable(False, False)

        hoje = date.today()
        frame = ttk.Frame(dialog, padding=10)
        frame.pack(fill="both", expand=True)

        ttk.Label(frame, text="Dia:").grid(row=0, column=0, sticky="w")
        dia_var = tk.StringVar(value=str(hoje.day))
        dia_cb = ttk.Combobox(frame, textvariable=dia_var, values=[str(i) for i in range(1,32)], width=4, state="readonly")
        dia_cb.grid(row=0, column=1, padx=(5, 15))

        ttk.Label(frame, text="Mês:").grid(row=0, column=2, sticky="w")
        meses_pt = [
            (1, "Jan"), (2, "Fev"), (3, "Mar"), (4, "Abr"), (5, "Mai"), (6, "Jun"),
            (7, "Jul"), (8, "Ago"), (9, "Set"), (10, "Out"), (11, "Nov"), (12, "Dez")
        ]
        mes_map_nome = {nome: num for num, nome in meses_pt}
        mes_var = tk.StringVar(value=meses_pt[hoje.month-1][1])
        mes_cb = ttk.Combobox(frame, textvariable=mes_var, values=[nome for _, nome in meses_pt], width=5, state="readonly")
        mes_cb.grid(row=0, column=3, padx=(5, 15))

        ttk.Label(frame, text="Ano:").grid(row=0, column=4, sticky="w")
        ano_var = tk.StringVar(value=str(hoje.year))
        anos_vals = [str(hoje.year + delta) for delta in range(-1, 3)]
        ano_cb = ttk.Combobox(frame, textvariable=ano_var, values=anos_vals, width=6, state="readonly")
        ano_cb.grid(row=0, column=5, padx=(5, 0))

        resultado = {"iso": None}

        def on_ok():
            try:
                d = int(dia_var.get())
                m = mes_map_nome.get(mes_var.get(), hoje.month)
                a = int(ano_var.get())
                dt = date(a, m, d)
                resultado["iso"] = dt.strftime("%Y-%m-%d")
                dialog.destroy()
            except Exception:
                messagebox.showerror("Data inválida", "Selecione uma data válida.", parent=dialog)

        def on_cancel():
            dialog.destroy()

        btns = ttk.Frame(frame)
        btns.grid(row=1, column=0, columnspan=6, pady=(12, 0))
        ttk.Button(btns, text="Cancelar", command=on_cancel).pack(side="right")
        ttk.Button(btns, text="OK", command=on_ok, style="Accent.TButton").pack(side="right", padx=(0, 8))

        dialog.bind("<Return>", lambda e: on_ok())
        dialog.bind("<Escape>", lambda e: on_cancel())
        dialog.wait_window()
        return resultado["iso"]

    def fetch_liturgia_por_data(self):
        data_str = self._abrir_dialogo_data()
        if not data_str:
            return

        self.status_label.config(text="Buscando liturgia online...")

        try:
            url = f"{LITURGIA_API_URL}?date={data_str}"
            response = requests.get(url, timeout=15)
            response.raise_for_status()
            data = response.json()

            readings = data.get("today", {}).get("readings", {})
            if not readings:
                raise ValueError("Nenhum dado de liturgia encontrado.")

            self._preencher_todas_secoes_de_readings(readings)

            self.status_label.config(text=f"Liturgia de {data_str} carregada com sucesso!")
            messagebox.showinfo("Sucesso", f"Liturgia de {data_str} importada. Verifique as seções atualizadas.", parent=self.master)

        except requests.RequestException as e:
            self.status_label.config(text="Erro ao buscar online.")
            messagebox.showerror("Erro de Rede", f"Falha ao conectar à API: {e}", parent=self.master)
        except Exception as e:
            self.status_label.config(text="Erro durante o processamento.")
            messagebox.showerror("Erro", f"Ocorreu um erro: {e}", parent=self.master)

    def _preencher_secao_leitura(self, nome_secao, dados_api):
        if not dados_api:
            return
        if nome_secao not in self.widgets_gui:
            return
        widgets = self.widgets_gui[nome_secao]

        # Seção de Aclamação (usa widgets específicos)
        if widgets.get("tipo") == "aclamacao":
            aclamacao = dados_api.get("head_response") or "Aleluia, Aleluia, Aleluia!"
            antifona = dados_api.get("head") or ""
            head_title = dados_api.get("head_title") or ""
            citacao = self._extrair_citacao_gospel(head_title)
            evangelista = self._extrair_evangelista(head_title)
            proclamacao = f"PROCLAMAÇÃO DO EVANGELHO DE {evangelista}" if evangelista else "PROCLAMAÇÃO DO EVANGELHO"
            linhas_antifona = []
            if antifona:
                linhas_antifona.append(antifona.strip())
            if citacao:
                linhas_antifona.append(citacao)

            widgets["aclamacao_widget"].delete("1.0", tk.END)
            widgets["aclamacao_widget"].insert(tk.END, "\n".join([proclamacao, aclamacao]))

            widgets["antifona_widget"].delete("1.0", tk.END)
            if linhas_antifona:
                widgets["antifona_widget"].insert(tk.END, "\n".join(linhas_antifona))
            return

        # Leitura comum (1ª e 2ª) e Salmo usam widgets de leitura
        if widgets.get("tipo") != "leitura":
            return

        if nome_secao == "Salmo":
            titulo = dados_api.get("title", "Salmo")
            resposta = dados_api.get("response", "")

            linhas_branco = []
            if resposta:
                linhas_branco.append(resposta)
  
            self._set_texto_leitura_widgets(widgets, titulo, linhas_branco)
            return

        # 1ª Leitura / 2ª Leitura: manter Título/Referência como padrão (PRIMEIRA/SEGUNDA LEITURA)
        # e colocar apenas a citação (ex.: "JUÍZES 2, 11-19") no TEXTO PRINCIPAL
        if nome_secao in ("1ª Leitura", "2ª Leitura"):
            titulo = dados_api.get("title", "")
            citacao = self._extrair_referencia_de_titulo(titulo) or ""
            try:
                aplicar_uppercase = True
                if "uppercase_var" in widgets:
                    try:
                        aplicar_uppercase = widgets["uppercase_var"].get()
                    except Exception:
                        aplicar_uppercase = True
                linhas = [citacao] if citacao else []
                if aplicar_uppercase:
                    linhas = [l.upper() for l in linhas]
                # Não mexe no titulo_amarelo_widget para manter PRIMEIRA/SEGUNDA LEITURA
                widgets["texto_branco_widget"].delete("1.0", tk.END)
                if linhas:
                    widgets["texto_branco_widget"].insert(tk.END, "\n".join(linhas))
            except Exception as e:
                print(f"Erro ao preencher citação da leitura: {e}")
            return

        # Outras leituras comuns (fallback)
        titulo = dados_api.get("title", "")
        texto = dados_api.get("text", "")
        ref = self._extrair_referencia_de_titulo(titulo)
        linhas_texto = [l.strip() for l in texto.split("\n") if l.strip()] if isinstance(texto, str) else []
        self._set_texto_leitura_widgets(widgets, ref or titulo, linhas_texto)

    def _preencher_todas_secoes_de_readings(self, readings):
        try:
            # 1ª Leitura
            self._preencher_secao_leitura("1ª Leitura", readings.get("first_reading", {}))
            # Salmo
            self._preencher_secao_leitura("Salmo", readings.get("psalm", {}))
            # 2ª Leitura
            segunda = readings.get("second_reading") or {}
            if segunda:
                self._preencher_secao_leitura("2ª Leitura", segunda)
            else:
                if "2ª Leitura" in self.widgets_gui:
                    w2 = self.widgets_gui["2ª Leitura"]
                    if w2.get("tipo") == "leitura":
                        try:
                            w2["titulo_amarelo_widget"].delete("1.0", tk.END)
                            w2["texto_branco_widget"].delete("1.0", tk.END)
                        except Exception:
                            pass
            # Aclamação
            self._preencher_secao_leitura("Aclamação", readings.get("gospel", {}))
        except Exception as e:
            print(f"Erro ao preencher seções: {e}")

    def carregar_liturgia_hoje_silencioso(self):
        try:
            data_str = date.today().strftime("%Y-%m-%d")
            url = f"{LITURGIA_API_URL}?date={data_str}"
            response = requests.get(url, timeout=15)
            response.raise_for_status()
            data = response.json()
            readings = data.get("today", {}).get("readings", {})
            if readings:
                self._preencher_todas_secoes_de_readings(readings)
        except Exception:
            # Silencioso: não mostrar mensagens
            pass

    def _set_texto_leitura_widgets(self, widgets, titulo_amarelo, linhas_texto):
        try:
            aplicar_uppercase = True
            if "uppercase_var" in widgets:
                try:
                    aplicar_uppercase = widgets["uppercase_var"].get()
                except Exception:
                    aplicar_uppercase = True
            titulo_final = titulo_amarelo or ""
            if aplicar_uppercase:
                titulo_final = titulo_final.upper()
                linhas_texto = [l.upper() for l in linhas_texto]

            widgets["titulo_amarelo_widget"].delete("1.0", tk.END)
            if titulo_final:
                widgets["titulo_amarelo_widget"].insert(tk.END, titulo_final)

            widgets["texto_branco_widget"].delete("1.0", tk.END)
            if linhas_texto:
                widgets["texto_branco_widget"].insert(tk.END, "\n".join(linhas_texto))
        except Exception as e:
            print(f"Erro ao preencher widgets de leitura: {e}")

    def _extrair_referencia_de_titulo(self, titulo):
        if not titulo:
            return ""
        # Remove prefixos como "Primeira leitura:" / "Segunda leitura:" mantendo a referência
        if ":" in titulo:
            return titulo.split(":", 1)[1].strip()
        return titulo.strip()

    def _extrair_citacao_gospel(self, head_title):
        if not head_title:
            return ""
        # Mapeia o evangelista para abreviação
        mapa = {
        "GÊNESIS": "Gn",
        "ÊXODO": "Ex",
        "LEVÍTICO": "Lv",
        "NÚMEROS": "Nm",
        "DEUTERONÔMIO": "Dt",
        "JOSUÉ": "Js",
        "JUÍZES": "Jz",
        "RUTE": "Rt",
        "1 SAMUEL": "1 Sm",
        "2 SAMUEL": "2 Sm",
        "1 REIS": "1 Rs",
        "2 REIS": "2 Rs",
        "1 CRÔNICAS": "1 Cr",
        "2 CRÔNICAS": "2 Cr",
        "ESDRAS": "Ed",
        "NEEMIAS": "Ne",
        "ESTER": "Est",
        "JÓ": "Jó",
        "SALMOS": "Sl",
        "PROVÉRBIOS": "Pv",
        "ECLESIASTES": "Ec",
        "CÂNTICO DOS CÂNTICOS": "Ct",
        "ISAÍAS": "Is",
        "JEREMIAS": "Jr",
        "LAMENTAÇÕES DE JEREMIAS": "Lm",
        "EZEQUIEL": "Ez",
        "DANIEL": "Dn",
        "OSÉIAS": "Os",
        "JOEL": "Jl",
        "AMÓS": "Am",
        "OBADIAS": "Ob",
        "JONAS": "Jn",
        "MIQUÉIAS": "Mq",
        "NAUM": "Na",
        "HABACUQUE": "Hc",
        "SOFONIAS": "Sf",
        "AGEU": "Ag",
        "ZACARIAS": "Zc",
        "MALAQUIAS": "Ml",
        "MATEUS": "Mt",
        "MARCOS": "Mc",
        "LUCAS": "Lc",
        "JOÃO": "Jo",
        "JOAO": "Jo", # Adicionado para flexibilidade na entrada
        "ATOS DOS APÓSTOLOS": "At",
        "ROMANOS": "Rm",
        "1 CORÍNTIOS": "1 Co",
        "2 CORÍNTIOS": "2 Co",
        "GÁLATAS": "Gl",
        "EFÉSIOS": "Ef",
        "FILIPENSES": "Fp",
        "COLOSSENSES": "Cl",
        "1 TESSALONICENSES": "1 Ts",
        "2 TESSALONICENSES": "2 Ts",
        "1 TIMÓTEO": "1 Tm",
        "2 TIMÓTEO": "2 Tm",
        "TITO": "Tt",
        "FILEMON": "Fm",
        "HEBREUS": "Hb",
        "TIAGO": "Tg",
        "1 PEDRO": "1 Pe",
        "2 PEDRO": "2 Pe",
        "1 JOÃO": "1 Jo",
        "2 JOÃO": "2 Jo",
        "3 JOÃO": "3 Jo",
        "JUDAS": "Jd",
        "APOCALIPSE": "Ap",
        "TOBIAS": "Tb",
        "JUDITE": "Jt",
        "BARUC": "Br",
        "ECLESIÁSTICO": "Eclo",
        "SABEDORIA": "Sb",
        "1 MACABEUS": "1 Mc",
        "2 MACABEUS": "2 Mc",
    }   
        ref_abrev = ""
        head_upper = head_title.upper()
        for nome, abrev in mapa.items():
            if nome in head_upper:
                ref_abrev = abrev
                break
        # Extrai padrão cap,versos (ex: 17, 14-20)
        m = re.search(r"(\d+\s*,\s*\d+(?:[\-–]\d+)?)", head_title)
        if not m:
            return ""
        nums = m.group(1)
        nums = nums.replace(" ", "").replace("–", "-")
        if ref_abrev:
            return f"({ref_abrev} {nums})"
        return f"({nums})"

    def _extrair_evangelista(self, head_title):
        if not head_title:
            return ""
        texto = head_title.upper()
        if "MATEUS" in texto:
            return "MATEUS"
        if "MARCOS" in texto:
            return "MARCOS"
        if "LUCAS" in texto:
            return "LUCAS"
        if "JOÃO" in texto or "JOAO" in texto:
            return "JOÃO"
        return ""






    def _atualizar_estado_botoes_controle_abas(self, event=None):
        try:
            if not self.notebook.tabs(): 
                if hasattr(self, 'remove_button'): self.remove_button.config(state="disabled")
                if hasattr(self, 'move_up_button'): self.move_up_button.config(state="disabled")
                if hasattr(self, 'move_down_button'): self.move_down_button.config(state="disabled")
                return

            selected_tab_id = self.notebook.select()
            if not selected_tab_id: 
                if hasattr(self, 'remove_button'): self.remove_button.config(state="normal" if self.notebook.tabs() else "disabled") 
                if hasattr(self, 'move_up_button'): self.move_up_button.config(state="disabled")
                if hasattr(self, 'move_down_button'): self.move_down_button.config(state="disabled")
                return

            if hasattr(self, 'remove_button'): self.remove_button.config(state="normal")
            current_index = self.notebook.index(selected_tab_id)
            total_tabs = self.notebook.index("end")

            if hasattr(self, 'move_up_button'): self.move_up_button.config(state="normal" if current_index > 0 else "disabled")
            if hasattr(self, 'move_down_button'): self.move_down_button.config(state="normal" if current_index < total_tabs - 1 else "disabled")

        except tk.TclError: 
            if hasattr(self, 'remove_button'): self.remove_button.config(state="disabled")
            if hasattr(self, 'move_up_button'): self.move_up_button.config(state="disabled")
            if hasattr(self, 'move_down_button'): self.move_down_button.config(state="disabled")


    def _reconstruir_ordem_geracao_dinamica(self):
        self.ordem_geracao_dinamica = []
        self.ordem_geracao_dinamica.append("TITULO_INICIAL_PLACEHOLDER")
        current_tabs_order = []
        if hasattr(self, 'notebook') and self.notebook.winfo_exists() and self.notebook.tabs(): 
            try: current_tabs_order = [self.notebook.tab(i, "text") for i in range(self.notebook.index("end"))]
            except tk.TclError: current_tabs_order = list(self.widgets_gui.keys()) if self.widgets_gui else self.ordem_gui_inicial
        else: current_tabs_order = list(self.widgets_gui.keys()) if self.widgets_gui else self.ordem_gui_inicial

        for tab_name in current_tabs_order:
            if tab_name in self.widgets_gui:
                self.ordem_geracao_dinamica.append(tab_name)
                tipo_secao_gui = self.widgets_gui[tab_name].get("tipo")
                nome_canonic_secao = self.DEFAULT_TEXTS.get(tab_name, {}).get("titulo", tab_name).upper()
                if tipo_secao_gui == "aclamacao" or nome_canonic_secao == "ACLAMAÇÃO DO EVANGELHO":
                    self.ordem_geracao_dinamica.extend(["CREDO", "PRECES"])
                elif tab_name == "Oferendas" or nome_canonic_secao == "PREPARAÇÃO DAS OFERENDAS":
                    self.ordem_geracao_dinamica.extend(["SANTO_TITULO", "ORACAO_EUCARISTICA", "CORDEIRO_TITULO"])
                elif tab_name == "Comunhão" or nome_canonic_secao == "COMUNHÃO":
                    self.ordem_geracao_dinamica.append("SANTA_LUZIA")
        
        if "Comunhão" not in current_tabs_order and "SANTA_LUZIA" not in self.ordem_geracao_dinamica and \
           any(s in self.ordem_geracao_dinamica for s in ["Oferendas", "Aclamação", "SANTO_TITULO"]):
            last_major_section_idx = -1
            potential_anchors = ["CORDEIRO_TITULO", "ORACAO_EUCARISTICA", "SANTO_TITULO", "Oferendas", "PRECES", "CREDO", "Aclamação"]
            for anchor in potential_anchors:
                if anchor in self.ordem_geracao_dinamica: last_major_section_idx = self.ordem_geracao_dinamica.index(anchor); break
            if last_major_section_idx != -1: self.ordem_geracao_dinamica.insert(last_major_section_idx + 1, "SANTA_LUZIA")
            elif current_tabs_order : 
                 last_tab_index_in_ordem = next((i for i, item in reversed(list(enumerate(self.ordem_geracao_dinamica))) if item in current_tabs_order), -1)
                 if last_tab_index_in_ordem != -1: self.ordem_geracao_dinamica.insert(last_tab_index_in_ordem + 1, "SANTA_LUZIA")
                 else: self.ordem_geracao_dinamica.append("SANTA_LUZIA")
            elif not self.ordem_geracao_dinamica or self.ordem_geracao_dinamica == ["TITULO_INICIAL_PLACEHOLDER"]:
                 self.ordem_geracao_dinamica.append("SANTA_LUZIA")

        if "AVISOS_IMG" in self.ordem_geracao_dinamica: self.ordem_geracao_dinamica.remove("AVISOS_IMG")
        self.ordem_geracao_dinamica.append("AVISOS_IMG")
        
        seen = set(); self.ordem_geracao_dinamica = [x for x in self.ordem_geracao_dinamica if not (x in seen or seen.add(x))]

    def _criar_aba_secao(self, nome_secao, tipo_override=None, inserir_em_posicao=-1, reconstruir_ordem=True):
        if nome_secao not in self.DEFAULT_TEXTS: 
             self.DEFAULT_TEXTS[nome_secao] = {
                "titulo": nome_secao.upper(), "refrao": [], "versos": [], "texto": [], "aclamacao_texto": [], 
                "antifona_texto": [], "titulo_amarelo": [], "texto_branco": []
            }
        frame = ttk.Frame(self.notebook, padding="10") 
        num_tabs_antes = self.notebook.index("end") if self.notebook.winfo_exists() else 0
        final_insert_pos = inserir_em_posicao
        if inserir_em_posicao == -1 or inserir_em_posicao >= num_tabs_antes: final_insert_pos = "end" 
        self.notebook.insert(final_insert_pos, frame, text=nome_secao)
        
        self.widgets_gui[nome_secao] = {}
        tipo_secao = tipo_override
        if not tipo_secao:
            if nome_secao in ["1ª Leitura", "Salmo", "2ª Leitura"]: tipo_secao = "leitura"
            elif nome_secao == "Aclamação": tipo_secao = "aclamacao"
            elif nome_secao == "Palavra": tipo_secao = "palavra"
            else: tipo_secao = "musica" 
        self.widgets_gui[nome_secao]["tipo"] = tipo_secao

        if tipo_secao == "leitura": self._criar_widgets_leitura_estilos(frame, nome_secao, self.widgets_gui[nome_secao])
        elif tipo_secao == "aclamacao": self._criar_widgets_aclamacao_estilos(frame, nome_secao, self.widgets_gui[nome_secao])
        elif tipo_secao == "palavra": self._criar_widgets_palavra_estilos(frame, nome_secao, self.widgets_gui[nome_secao])
        else: self._criar_widgets_musica_estilos(frame, nome_secao, self.widgets_gui[nome_secao])
        
        if reconstruir_ordem: self._reconstruir_ordem_geracao_dinamica()
        self._atualizar_estado_botoes_controle_abas() 

    def dialogo_adicionar_secao(self):
        dialog = tk.Toplevel(self.master); dialog.title("Adicionar Nova Seção Musical")
        dialog.geometry("380x180"); dialog.transient(self.master); dialog.grab_set(); dialog.resizable(False, False)
        main_frame = ttk.Frame(dialog, padding=10); main_frame.pack(expand=True, fill="both")
        ttk.Label(main_frame, text="Nome da Nova Seção:").grid(row=0, column=0, sticky="w", pady=(0,2))
        nome_entry = ttk.Entry(main_frame, width=40); nome_entry.grid(row=1, column=0, columnspan=2, sticky="ew", pady=(0,10)); nome_entry.focus_set()
        ttk.Label(main_frame, text="Inserir após:").grid(row=2, column=0, sticky="w", pady=(0,2))
        posicoes_disponiveis = ["No início"]
        try:
            num_tabs = self.notebook.index("end")
            if num_tabs > 0: posicoes_disponiveis.extend([self.notebook.tab(i, "text") for i in range(num_tabs)])
            posicoes_disponiveis.append("No fim")
        except tk.TclError: posicoes_disponiveis = ["No início", "No fim"]
        posicao_var = tk.StringVar(value=posicoes_disponiveis[-1]) 
        posicao_combo = ttk.Combobox(main_frame, textvariable=posicao_var, values=posicoes_disponiveis, state="readonly", width=37); posicao_combo.grid(row=3, column=0, columnspan=2, sticky="ew", pady=(0,15))
        def on_ok():
            nome_nova_secao = nome_entry.get().strip()
            if not nome_nova_secao: messagebox.showerror("Erro", "O nome da seção não pode ser vazio.", parent=dialog); return
            if nome_nova_secao in self.widgets_gui: messagebox.showerror("Erro", f"A seção '{nome_nova_secao}' já existe.", parent=dialog); return
            posicao_selecionada = posicao_var.get(); idx_insercao = -1 
            if self.notebook.winfo_exists() and self.notebook.index("end") is not None: idx_insercao = self.notebook.index("end")
            else: idx_insercao = 0
            if posicao_selecionada == "No início": idx_insercao = 0
            elif posicao_selecionada != "No fim": 
                try:
                    abas_existentes = [self.notebook.tab(i, "text") for i in range(self.notebook.index("end"))]
                    if posicao_selecionada in abas_existentes: idx_insercao = abas_existentes.index(posicao_selecionada) + 1
                except tk.TclError: pass 
            if nome_nova_secao not in self.DEFAULT_TEXTS: self.DEFAULT_TEXTS[nome_nova_secao] = {"titulo": nome_nova_secao.upper(), "refrao": [], "versos": []} 
            self._criar_aba_secao(nome_nova_secao, tipo_override="musica", inserir_em_posicao=idx_insercao, reconstruir_ordem=True)
            dialog.destroy()
        button_frame = ttk.Frame(main_frame); button_frame.grid(row=4, column=0, columnspan=2, pady=(5,0))
        ok_button = ttk.Button(button_frame, text="Adicionar", command=on_ok, style="Accent.TButton"); ok_button.pack(side="right", padx=(5,0))
        cancel_button = ttk.Button(button_frame, text="Cancelar", command=dialog.destroy); cancel_button.pack(side="right")
        dialog.bind("<Return>", lambda event: on_ok()); dialog.bind("<Escape>", lambda event: dialog.destroy())

    def remover_secao_selecionada(self):
        if not self.notebook.tabs(): messagebox.showinfo("Remover Seção", "Nenhuma seção para remover.", parent=self.master); return
        try:
            selected_tab_id = self.notebook.select()
            if not selected_tab_id: messagebox.showinfo("Remover Seção", "Nenhuma seção selecionada.", parent=self.master); return
            tab_name = self.notebook.tab(selected_tab_id, "text")
        except tk.TclError: messagebox.showinfo("Remover Seção", "Nenhuma seção selecionada para remover.", parent=self.master); return
        confirm_msg = f"Tem certeza que deseja remover a seção '{tab_name}'?"
        if messagebox.askyesno("Confirmar Remoção", confirm_msg, parent=self.master):
            self.notebook.forget(selected_tab_id)
            if tab_name in self.widgets_gui: del self.widgets_gui[tab_name]
            if tab_name in self.DEFAULT_TEXTS: del self.DEFAULT_TEXTS[tab_name] 
            self._reconstruir_ordem_geracao_dinamica()
            self._atualizar_estado_botoes_controle_abas() 
            self.status_label.config(text=f"Seção '{tab_name}' removida.")

    def mover_secao_selecionada(self, direcao): 
        if not self.notebook.tabs(): return
        try:
            selected_tab_id = self.notebook.select()
            if not selected_tab_id: return
            current_index = self.notebook.index(selected_tab_id)
            new_index = current_index + direcao

            if 0 <= new_index < self.notebook.index("end"):
                self.notebook.insert(new_index, selected_tab_id) 
                self.notebook.select(new_index) 
                self._reconstruir_ordem_geracao_dinamica()
                self._atualizar_estado_botoes_controle_abas()
                tab_name = self.notebook.tab(new_index, "text")
                self.status_label.config(text=f"Seção '{tab_name}' movida.")
        except tk.TclError as e:
            print(f"Erro ao mover aba: {e}")
            pass 

    def _criar_controles_estilo(self, parent, data_dict, prefix_key, default_size_pt, default_bold=True, default_italic=False):
        style_frame = ttk.Frame(parent); style_frame.pack(fill='x', expand=True, pady=2)
        controls_left_frame = ttk.Frame(style_frame)
        controls_left_frame.pack(side="left", fill="x", expand=True, padx=(0,10))
        font_label = ttk.Label(controls_left_frame, text="Fonte:")
        font_label.grid(row=0, column=0, sticky='w', padx=(0, 2), pady=2)
        font_var = tk.StringVar(value=NOME_FONTE_PADRAO)
        font_combo = ttk.Combobox(controls_left_frame, textvariable=font_var, values=FONTES_COMUNS_PPT, width=23, state='readonly', style="TCombobox")
        font_combo.grid(row=0, column=1, sticky='ew', padx=(0, 10), pady=2)
        data_dict[f"{prefix_key}_font_combo"] = font_combo
        size_label = ttk.Label(controls_left_frame, text="Tam:")
        size_label.grid(row=0, column=2, sticky='w', padx=(0, 2), pady=2)
        size_spinbox = ttk.Spinbox(controls_left_frame, from_=10, to=120, increment=1, width=5, justify='right', wrap=True, style="TSpinbox")
        size_spinbox.set(int(default_size_pt.pt))
        size_spinbox.grid(row=0, column=3, sticky='w', pady=2)
        data_dict[f"{prefix_key}_font_spinbox"] = size_spinbox
        bold_var = tk.BooleanVar(value=default_bold)
        bold_check = ttk.Checkbutton(controls_left_frame, text="Negrito", variable=bold_var, style="TCheckbutton")
        bold_check.grid(row=1, column=0, columnspan=2, sticky='w', pady=2)
        data_dict[f"{prefix_key}_bold_var"] = bold_var
        italic_var = tk.BooleanVar(value=default_italic)
        italic_check = ttk.Checkbutton(controls_left_frame, text="Itálico", variable=italic_var, style="TCheckbutton")
        italic_check.grid(row=1, column=2, columnspan=2, sticky='w', padx=(10,0), pady=2) 
        data_dict[f"{prefix_key}_italic_var"] = italic_var
        controls_left_frame.columnconfigure(1, weight=1) 
        preview_label = tk.Label(style_frame, text="Amostra", font=(NOME_FONTE_PADRAO, 11), width=12, height=2, relief="groove", borderwidth=1, bg='white') 
        preview_label.pack(side="right", fill="y", padx=(5,0))
        data_dict[f"{prefix_key}_preview_label"] = preview_label
        def update_preview(*args):
            try: 
                fname = font_var.get(); fsize_val = size_spinbox.get()
                if not fsize_val: fsize = 12 
                else: fsize = int(fsize_val)
                fbold = "bold" if bold_var.get() else "normal"; fitalic = "italic" if italic_var.get() else "roman"
                display_fsize = fsize if fsize < 16 else 16 
                preview_label.config(font=(fname, display_fsize, fbold, fitalic), text=fname)
            except (ValueError, tk.TclError): preview_label.config(font=(NOME_FONTE_PADRAO, 11), text="?")
        font_var.trace_add("write", update_preview); size_spinbox.config(command=update_preview) 
        bold_var.trace_add("write", update_preview); italic_var.trace_add("write", update_preview)
        update_preview()

    def _criar_widgets_musica_estilos(self, parent_frame, nome_parte, data_dict):
        title_config_frame = ttk.Frame(parent_frame); title_config_frame.pack(fill='x', pady=(0,10)) 
        ttk.Label(title_config_frame, text="Título da Seção:", font=('Arial', 10, 'bold')).pack(side='left', anchor='w', padx=(0,5))
        default_titulo_secao = self.DEFAULT_TEXTS.get(nome_parte, {}).get("titulo", nome_parte.upper())
        data_dict["titulo_secao_entry"] = ttk.Entry(title_config_frame, width=60, font=('Arial', 10))
        data_dict["titulo_secao_entry"].insert(0, default_titulo_secao)
        data_dict["titulo_secao_entry"].pack(side='left', fill='x', expand=True)
        options_frame = ttk.Frame(parent_frame) 
        options_frame.pack(fill='x', pady=(0,5))
        data_dict["iniciar_refrao_var"] = tk.BooleanVar(value=False); chk_start_refrao = ttk.Checkbutton(options_frame, text="Iniciar com Refrão", variable=data_dict["iniciar_refrao_var"]); chk_start_refrao.pack(side='left', anchor='w', padx=(0,15))
        data_dict["uppercase_var"] = tk.BooleanVar(value=True)
        chk_uppercase = ttk.Checkbutton(options_frame, text="Texto em Maiúsculas", variable=data_dict["uppercase_var"])
        chk_uppercase.pack(side='left', anchor='w')
        # Opções de detecção automática de refrão
        data_dict["detectar_refrao_auto_var"] = tk.BooleanVar(value=True)
        chk_detect = ttk.Checkbutton(options_frame, text="Detectar e organizar refrão automaticamente", variable=data_dict["detectar_refrao_auto_var"]) 
        chk_detect.pack(side='left', anchor='w', padx=(15,0))
        # Ao ativar, se já houver texto, processa imediatamente
        def _on_toggle_detect(*_):
            try:
                if data_dict["detectar_refrao_auto_var"].get():
                    self._processar_refrao_auto_se_existe(nome_parte)
            except Exception:
                pass
        try:
            data_dict["detectar_refrao_auto_var"].trace_add('write', _on_toggle_detect)
        except Exception:
            pass
        refrao_frame = ttk.LabelFrame(parent_frame, text="Refrão (Amarelo)", padding=10); refrao_frame.pack(fill='x', expand=False, pady=5)
        self._criar_controles_estilo(refrao_frame, data_dict, "refrao", DEFAULT_TAMANHO_FONTE_MUSICA_REFRAO, default_bold=True)
        data_dict["refrao_widget"] = scrolledtext.ScrolledText(refrao_frame, height=6, width=90, wrap=tk.WORD, font=('Arial', 10)); data_dict["refrao_widget"].pack(fill="x", expand=True, padx=0, pady=(5,0))
        default_refrao = self.DEFAULT_TEXTS.get(nome_parte, {}).get("refrao", []) 
        if default_refrao: data_dict["refrao_widget"].insert(tk.END, "\n".join(default_refrao))
        # Bind: se usuário colar música inteira no Refrão, reorganiza (mantém só refrão e move versos)
        def _schedule_from_refrao(event=None, nome=nome_parte):
            try:
                if not data_dict.get("detectar_refrao_auto_var").get():
                    return
            except Exception:
                return
            self.master.after(120, lambda: self._reorganizar_quando_colado_no_refrao(nome))
        try:
            data_dict["refrao_widget"].bind('<<Paste>>', _schedule_from_refrao)
            data_dict["refrao_widget"].bind('<Control-v>', _schedule_from_refrao)
            data_dict["refrao_widget"].bind('<Control-V>', _schedule_from_refrao)
        except Exception:
            pass
        verso_frame = ttk.LabelFrame(parent_frame, text="Versos (Branco)", padding=10); verso_frame.pack(fill='both', expand=True, pady=5)
        self._criar_controles_estilo(verso_frame, data_dict, "verso", DEFAULT_TAMANHO_FONTE_MUSICA_VERSO, default_bold=True)
        data_dict["verso_widget"] = scrolledtext.ScrolledText(verso_frame, height=10, width=90, wrap=tk.WORD, font=('Arial', 10)); data_dict["verso_widget"].pack(fill="both", expand=True, padx=0, pady=(5,0))
        default_versos_list_of_lists = self.DEFAULT_TEXTS.get(nome_parte, {}).get("versos", []) 
        default_versos_text = ["\n".join(estrofe) for estrofe in default_versos_list_of_lists]
        if default_versos_text: data_dict["verso_widget"].insert(tk.END, "\n\n".join(default_versos_text))
        # Bind: detectar refrão após colar
        def _schedule_detect_refrao_event(event=None, nome=nome_parte):
            try:
                if not data_dict.get("detectar_refrao_auto_var").get():
                    return
            except Exception:
                return
            self.master.after(120, lambda: self._detectar_e_aplicar_refrao(nome))
        try:
            data_dict["verso_widget"].bind('<<Paste>>', _schedule_detect_refrao_event)
            data_dict["verso_widget"].bind('<Control-v>', _schedule_detect_refrao_event)
            data_dict["verso_widget"].bind('<Control-V>', _schedule_detect_refrao_event)
        except Exception:
            pass

    def _criar_widgets_leitura_estilos(self, parent_frame, nome_parte, data_dict):
        titulo_amarelo_padrao_list = self.DEFAULT_TEXTS.get(nome_parte, {}).get("titulo_amarelo", [nome_parte.upper()])
        ref_frame = ttk.LabelFrame(parent_frame, text=f"Título/Referência - {nome_parte} (Amarelo)", padding=10); ref_frame.pack(fill='x', expand=False, pady=5)
        self._criar_controles_estilo(ref_frame, data_dict, "titulo_amarelo", DEFAULT_TAMANHO_FONTE_LEITURA_TITULO_AMARELO, default_bold=True)
        data_dict["titulo_amarelo_widget"] = scrolledtext.ScrolledText(ref_frame, height=4, width=90, wrap=tk.WORD, font=('Arial', 10)); data_dict["titulo_amarelo_widget"].pack(fill="x", expand=True, padx=0, pady=(5,0))
        if titulo_amarelo_padrao_list: data_dict["titulo_amarelo_widget"].insert(tk.END, "\n".join(titulo_amarelo_padrao_list))
        texto_frame = ttk.LabelFrame(parent_frame, text=f"Texto Principal - {nome_parte} (Branco)", padding=10); texto_frame.pack(fill='both', expand=True, pady=5)
        self._criar_controles_estilo(texto_frame, data_dict, "texto_branco", DEFAULT_TAMANHO_FONTE_LEITURA_TEXTO_BRANCO, default_bold=True)
        data_dict["texto_branco_widget"] = scrolledtext.ScrolledText(texto_frame, height=15, width=90, wrap=tk.WORD, font=('Arial', 10)); data_dict["texto_branco_widget"].pack(fill="both", expand=True, padx=0, pady=(5,0))
        default_txt = self.DEFAULT_TEXTS.get(nome_parte, {}).get("texto_branco", [])
        if default_txt: data_dict["texto_branco_widget"].insert(tk.END, "\n".join(default_txt))

    def _detectar_e_aplicar_refrao(self, nome_parte: str):
        try:
            if nome_parte not in self.widgets_gui:
                return
            data_gui = self.widgets_gui[nome_parte]
            if data_gui.get("tipo") != "musica":
                return
            versos_texto = data_gui["verso_widget"].get("1.0", tk.END).strip()
            if not versos_texto:
                return
            # Divide por estrofes (blocos) separados por linha em branco
            estrofes_raw = [b for b in versos_texto.split("\n\n")]
            estrofes = []
            for bloco in estrofes_raw:
                linhas = [l.strip() for l in bloco.split("\n") if l.strip()]
                if linhas:
                    estrofes.append(linhas)
            if not estrofes:
                return
            def _norm(block):
                return "\n".join([ln.lower().strip() for ln in block if ln.strip()])
            # Seleciona o refrão pela regra: maior frequência; empate → primeiro que aparece
            idx_map = {}
            freq = {}
            for i, blk in enumerate(estrofes):
                k = _norm(blk)
                if not k:
                    continue
                freq[k] = freq.get(k, 0) + 1
                idx_map.setdefault(k, i)
            # Encontrar melhor candidato (apenas se aparecer 2+ vezes)
            candidatos = [(count, -idx_map[k], k) for k, count in freq.items() if count >= 2]
            if not candidatos:
                return
            candidatos.sort(reverse=True)
            chave_refrao = candidatos[0][2]
            # Construir novos versos sem o refrão escolhido (remove todas as ocorrências do candidato)
            estrofes_filtradas = [blk for blk in estrofes if _norm(blk) != chave_refrao]
            # Refrão a partir da primeira ocorrência do candidato
            exemplar = next((blk for blk in estrofes if _norm(blk) == chave_refrao), None)
            novos_refroes = ["\n".join(exemplar)] if exemplar else []
            # Atualiza widgets
            # Refrão: anexa mantendo separação dupla entre blocos
            refrao_atual = data_gui["refrao_widget"].get("1.0", tk.END).strip()
            refrao_blocos = [refrao_atual] if refrao_atual else []
            for r in novos_refroes:
                if r and (not refrao_atual or r.lower() not in refrao_atual.lower()):
                    refrao_blocos.append(r)
            data_gui["refrao_widget"].delete("1.0", tk.END)
            if refrao_blocos:
                data_gui["refrao_widget"].insert(tk.END, "\n\n".join([b for b in refrao_blocos if b]))
            # Versos restantes
            data_gui["verso_widget"].delete("1.0", tk.END)
            if estrofes_filtradas:
                texto_restante = "\n\n".join(["\n".join(blk) for blk in estrofes_filtradas])
                data_gui["verso_widget"].insert(tk.END, texto_restante)
            self.status_label.config(text=f"Refrão detectado e movido em '{nome_parte}'.")
        except Exception as e:
            print(f"Falha na detecção automática de refrão: {e}")

    def _reorganizar_quando_colado_no_refrao(self, nome_parte: str):
        try:
            if nome_parte not in self.widgets_gui:
                return
            data_gui = self.widgets_gui[nome_parte]
            if data_gui.get("tipo") != "musica":
                return
            texto_refrao = data_gui["refrao_widget"].get("1.0", tk.END).strip()
            if not texto_refrao:
                return
            # Divide por estrofes
            blocos_raw = [b for b in texto_refrao.split("\n\n")]
            blocos = []
            for bloco in blocos_raw:
                linhas = [l.strip() for l in bloco.split("\n") if l.strip()]
                if linhas:
                    blocos.append(linhas)
            if not blocos:
                return
            def _norm(block):
                return "\n".join([ln.lower().strip() for ln in block if ln.strip()])
            idx_map = {}
            freq = {}
            for i, blk in enumerate(blocos):
                k = _norm(blk)
                if not k:
                    continue
                freq[k] = freq.get(k, 0) + 1
                idx_map.setdefault(k, i)
            candidatos = [(count, -idx_map[k], k) for k, count in freq.items()]
            candidatos.sort(reverse=True)
            chave_refrao = candidatos[0][2] if candidatos else _norm(blocos[0])
            # Refrão: apenas o bloco candidato
            exemplar = next((blk for blk in blocos if _norm(blk) == chave_refrao), None)
            refr_set = ["\n".join(exemplar)] if exemplar else ["\n".join(blocos[0])]
            # Versos restantes: remove todas as ocorrências do candidato
            versos_restantes = ["\n".join(blk) for blk in blocos if _norm(blk) != chave_refrao]
            # Atualiza widgets
            data_gui["refrao_widget"].delete("1.0", tk.END)
            if refr_set:
                data_gui["refrao_widget"].insert(tk.END, "\n\n".join(refr_set))
            data_gui["verso_widget"].delete("1.0", tk.END)
            if versos_restantes:
                data_gui["verso_widget"].insert(tk.END, "\n\n".join(versos_restantes))
            self.status_label.config(text=f"Refrão/Versos reorganizados em '{nome_parte}'.")
        except Exception as e:
            print(f"Falha ao reorganizar colagem no refrão: {e}")

    def _processar_refrao_auto_se_existe(self, nome_parte: str):
        # Se o checkbox auto estiver ligado e já houver conteúdo, decide qual rotina usar
        try:
            if nome_parte not in self.widgets_gui:
                return
            data_gui = self.widgets_gui[nome_parte]
            if data_gui.get("tipo") != "musica":
                return
            versos_texto = data_gui["verso_widget"].get("1.0", tk.END).strip()
            refrao_texto = data_gui["refrao_widget"].get("1.0", tk.END).strip()
            if versos_texto and not refrao_texto:
                self._detectar_e_aplicar_refrao(nome_parte)
            elif refrao_texto and not versos_texto:
                self._reorganizar_quando_colado_no_refrao(nome_parte)
            elif versos_texto and refrao_texto:
                # Se ambos possuem conteúdo, prioriza consolidar refrão (repetições) a partir dos versos
                self._detectar_e_aplicar_refrao(nome_parte)
        except Exception:
            pass

    def _criar_widgets_aclamacao_estilos(self, parent_frame, nome_parte, data_dict):
        title_config_frame = ttk.Frame(parent_frame); title_config_frame.pack(fill='x', pady=(0,5))
        ttk.Label(title_config_frame, text="Título da Seção:", font=('Arial', 10, 'bold')).pack(side='left', anchor='w', padx=(0,5))
        default_titulo_secao = self.DEFAULT_TEXTS.get(nome_parte, {}).get("titulo", nome_parte.upper())
        data_dict["titulo_secao_entry"] = ttk.Entry(title_config_frame, width=60, font=('Arial', 10))
        data_dict["titulo_secao_entry"].insert(0, default_titulo_secao)
        data_dict["titulo_secao_entry"].pack(side='left', fill='x', expand=True)
        uppercase_options_frame = ttk.Frame(parent_frame); uppercase_options_frame.pack(fill='x', pady=(0,10)) 
        data_dict["uppercase_var"] = tk.BooleanVar(value=True)
        chk_uppercase = ttk.Checkbutton(uppercase_options_frame, text="Textos em Maiúsculas", variable=data_dict["uppercase_var"])
        chk_uppercase.pack(side='left', anchor='w')
        aclamacao_frame = ttk.LabelFrame(parent_frame, text="Aclamação (Amarelo - Superior)", padding=10); aclamacao_frame.pack(fill='x', expand=False, pady=5)
        self._criar_controles_estilo(aclamacao_frame, data_dict, "aclamacao", DEFAULT_TAMANHO_FONTE_ACLAMACAO, default_bold=True)
        data_dict["aclamacao_widget"] = scrolledtext.ScrolledText(aclamacao_frame, height=5, width=90, wrap=tk.WORD, font=('Arial', 10)); data_dict["aclamacao_widget"].pack(fill="x", expand=True, padx=0, pady=(5,0))
        default_aclamacao = self.DEFAULT_TEXTS.get(nome_parte, {}).get("aclamacao_texto", []);
        if default_aclamacao: data_dict["aclamacao_widget"].insert(tk.END, "\n".join(default_aclamacao))
        antifona_frame = ttk.LabelFrame(parent_frame, text="Antífona (Branco - Inferior)", padding=10); antifona_frame.pack(fill='both', expand=True, pady=5)
        self._criar_controles_estilo(antifona_frame, data_dict, "antifona", DEFAULT_TAMANHO_FONTE_ANTIFONA, default_bold=True)
        data_dict["antifona_widget"] = scrolledtext.ScrolledText(antifona_frame, height=12, width=90, wrap=tk.WORD, font=('Arial', 10)); data_dict["antifona_widget"].pack(fill="both", expand=True, padx=0, pady=(5,0))
        default_antifona = self.DEFAULT_TEXTS.get(nome_parte, {}).get("antifona_texto", []);
        if default_antifona: data_dict["antifona_widget"].insert(tk.END, "\n".join(default_antifona))

    def _criar_widgets_palavra_estilos(self, parent_frame, nome_parte, data_dict):
        title_config_frame = ttk.Frame(parent_frame); title_config_frame.pack(fill='x', pady=(0,5))
        ttk.Label(title_config_frame, text="Título da Seção:", font=('Arial', 10, 'bold')).pack(side='left', anchor='w', padx=(0,5))
        default_titulo_secao = self.DEFAULT_TEXTS.get(nome_parte, {}).get("titulo", nome_parte.upper())
        data_dict["titulo_secao_entry"] = ttk.Entry(title_config_frame, width=60, font=('Arial', 10))
        data_dict["titulo_secao_entry"].insert(0, default_titulo_secao)
        data_dict["titulo_secao_entry"].pack(side='left', fill='x', expand=True)
        uppercase_options_frame = ttk.Frame(parent_frame); uppercase_options_frame.pack(fill='x', pady=(0,10))
        data_dict["uppercase_var"] = tk.BooleanVar(value=True)
        chk_uppercase = ttk.Checkbutton(uppercase_options_frame, text="Texto em Maiúsculas", variable=data_dict["uppercase_var"])
        chk_uppercase.pack(side='left', anchor='w')
        texto_frame = ttk.LabelFrame(parent_frame, text=f"Texto - {nome_parte} (Amarelo)", padding=10); texto_frame.pack(fill='both', expand=True, pady=5)
        self._criar_controles_estilo(texto_frame, data_dict, "texto", DEFAULT_TAMANHO_FONTE_PALAVRA, default_bold=True)
        data_dict["texto_widget"] = scrolledtext.ScrolledText(texto_frame, height=20, width=90, wrap=tk.WORD, font=('Arial', 10)); data_dict["texto_widget"].pack(fill="both", expand=True, padx=0, pady=(5,0))
        default_txt = self.DEFAULT_TEXTS.get(nome_parte, {}).get("texto", [])
        if default_txt: data_dict["texto_widget"].insert(tk.END, "\n".join(default_txt))

    def configurar_oracoes_eucaristicas_dialog(self):
        dialog = tk.Toplevel(self.master)
        dialog.title("Selecionar Orações Eucarísticas")
        dialog.geometry("650x560")
        dialog.transient(self.master)
        dialog.grab_set()
        dialog.resizable(True, True)

        main = ttk.Frame(dialog, padding=10)
        main.pack(fill="both", expand=True)

        top_bar = ttk.Frame(main)
        top_bar.pack(fill="x")
        chk_all = ttk.Checkbutton(top_bar, text="Selecionar todas", variable=self.eucaristica_selecionar_todas, command=lambda: self._toggle_todas_eucaristicas(listbox))
        chk_all.pack(side="left")
        # Fonte
        ttk.Label(top_bar, text="Tam. fonte:").pack(side="left", padx=(20,4))
        font_var = tk.StringVar(value=str(self.eucaristica_font_size_pt))
        font_spin = ttk.Spinbox(top_bar, from_=40, to=120, increment=2, width=5, textvariable=font_var)
        font_spin.pack(side="left")

        container = ttk.Frame(main)
        container.pack(fill="both", expand=True, pady=(10,5))

        listbox = tk.Listbox(container, selectmode=tk.MULTIPLE)
        scrollbar = ttk.Scrollbar(container, orient="vertical", command=listbox.yview)
        listbox.config(yscrollcommand=scrollbar.set)
        listbox.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")

        nomes = list(ORACOES_EUCARISTICAS.keys())
        for nome in nomes:
            listbox.insert(tk.END, nome)

        # Seleção inicial: se "selecionar todas" estiver marcado ao abrir, refletir na listbox
        if self.eucaristica_selecionar_todas.get():
            try:
                listbox.select_set(0, tk.END)
            except Exception:
                pass

        preview = scrolledtext.ScrolledText(main, height=10)
        preview.pack(fill="both", expand=False, pady=(10,5))

        def atualizar_preview(event=None):
            preview.delete("1.0", tk.END)
            selecao = [listbox.get(i) for i in listbox.curselection()]
            if not selecao:
                return
            texto = []
            for nome in selecao:
                texto.append(nome)
                texto.extend(ORACOES_EUCARISTICAS.get(nome, []))
                texto.append("\n")
            preview.insert(tk.END, "\n".join(texto))

        listbox.bind("<<ListboxSelect>>", atualizar_preview)

        btns = ttk.Frame(main)
        btns.pack(fill="x")
        def aplicar():
            self.eucaristica_escolhidas = set([listbox.get(i) for i in listbox.curselection()])
            try:
                val = int(font_var.get())
                if 10 <= val <= 160:
                    self.eucaristica_font_size_pt = val
            except Exception:
                pass
            dialog.destroy()
        def cancelar():
            dialog.destroy()
        ttk.Button(btns, text="Cancelar", command=cancelar).pack(side="right")
        ttk.Button(btns, text="Aplicar", command=aplicar, style="Accent.TButton").pack(side="right", padx=(0,8))

    def _toggle_todas_eucaristicas(self, listbox):
        try:
            if self.eucaristica_selecionar_todas.get():
                listbox.select_set(0, tk.END)
            else:
                listbox.selection_clear(0, tk.END)
        except Exception:
            pass

    def _obter_oracoes_eucaristicas_para_geracao(self):
        if self.eucaristica_selecionar_todas.get() or not self.eucaristica_escolhidas:
            return [ORACOES_EUCARISTICAS[n] for n in ORACOES_EUCARISTICAS.keys()]
        return [ORACOES_EUCARISTICAS[n] for n in ORACOES_EUCARISTICAS.keys() if n in self.eucaristica_escolhidas]

    def _obter_oracoes_eucaristicas_com_nomes(self):
        # Retorna lista de tuplas (nome, bloco)
        if self.eucaristica_selecionar_todas.get() or not self.eucaristica_escolhidas:
            return [(nome, ORACOES_EUCARISTICAS[nome]) for nome in ORACOES_EUCARISTICAS.keys()]
        else:
            return [(nome, ORACOES_EUCARISTICAS[nome]) for nome in ORACOES_EUCARISTICAS.keys() if nome in self.eucaristica_escolhidas]

    def _get_font_style_from_gui(self,data_dict, prefix_key, default_size_pt, default_bold=True, default_italic=False):
        font_size = default_size_pt; font_name = NOME_FONTE_PADRAO; bold_state = default_bold; italic_state = default_italic
        spinbox = data_dict.get(f"{prefix_key}_font_spinbox"); combo = data_dict.get(f"{prefix_key}_font_combo"); bold_var = data_dict.get(f"{prefix_key}_bold_var"); italic_var = data_dict.get(f"{prefix_key}_italic_var")
        if spinbox:
            try: valor_int = int(spinbox.get()); font_size = Pt(valor_int) if 10 <= valor_int <= 120 else default_size_pt
            except (tk.TclError, ValueError): pass
        if combo: selected_font = combo.get(); font_name = selected_font if selected_font in FONTES_COMUNS_PPT else NOME_FONTE_PADRAO
        if bold_var:
            try: bold_state = bold_var.get()
            except tk.TclError: pass
        if italic_var:
            try: italic_state = italic_var.get()
            except tk.TclError: pass
        return font_size, font_name, bold_state, italic_state

    def adicionar_secao_musical(self, prs, layout_slide_branco, nome_parte_gui):
        conteudo_adicionado_total = False; data_gui = self.widgets_gui.get(nome_parte_gui)
        pegar_da_gui = bool(data_gui) 
        titulo_parte = self.DEFAULT_TEXTS.get(nome_parte_gui, {}).get("titulo", nome_parte_gui.upper())
        if pegar_da_gui and "titulo_secao_entry" in data_gui:
            gui_titulo = data_gui["titulo_secao_entry"].get().strip();
            if gui_titulo: titulo_parte = gui_titulo
        refrao_gui_str = ""; versos_gui_str = ""; tamanho_refrao, nome_refrao, bold_refrao, italic_refrao = DEFAULT_TAMANHO_FONTE_MUSICA_REFRAO, NOME_FONTE_PADRAO, True, False
        tamanho_verso, nome_verso, bold_verso, italic_verso = DEFAULT_TAMANHO_FONTE_MUSICA_VERSO, NOME_FONTE_PADRAO, True, False
        iniciar_com_refrao = False; aplicar_uppercase = True ; refrao_final = []; versos_processados = []
        if pegar_da_gui:
            refrao_gui_str = data_gui["refrao_widget"].get("1.0", tk.END).strip(); versos_gui_str = data_gui["verso_widget"].get("1.0", tk.END).strip()
            tamanho_refrao, nome_refrao, bold_refrao, italic_refrao = self._get_font_style_from_gui(data_gui, "refrao", DEFAULT_TAMANHO_FONTE_MUSICA_REFRAO, True, False)
            tamanho_verso, nome_verso, bold_verso, italic_verso = self._get_font_style_from_gui(data_gui, "verso", DEFAULT_TAMANHO_FONTE_MUSICA_VERSO, True, False)
            if "iniciar_refrao_var" in data_gui:
                try: iniciar_com_refrao = data_gui["iniciar_refrao_var"].get()
                except Exception: iniciar_com_refrao = False
            if "uppercase_var" in data_gui:
                try: aplicar_uppercase = data_gui["uppercase_var"].get()
                except Exception: aplicar_uppercase = True
            refrao_final = [l.strip() for l in refrao_gui_str.split('\n') if l.strip()] if refrao_gui_str else []
            if versos_gui_str:
                estrofes_input = versos_gui_str.split('\n\n')
                for estrofe_str in estrofes_input:
                    if estrofe_str.strip():
                        linhas_estrofe = [l.strip() for l in estrofe_str.split('\n') if l.strip()]
                        if linhas_estrofe: versos_processados.append(linhas_estrofe)
            # Detectar estrofes repetidas como refrão e mover
            try:
                if versos_processados:
                    # Normalização para comparação
                    def _norm(block):
                        return '\n'.join([ln.lower().strip() for ln in block if ln.strip()])
                    freq = {}
                    for bloco in versos_processados:
                        key = _norm(bloco)
                        if not key:
                            continue
                        freq[key] = freq.get(key, 0) + 1
                    # Quais são refrões (aparecem 2+ vezes)
                    chaves_refrao = {k for k, c in freq.items() if c >= 2}
                    if chaves_refrao:
                        # Adiciona ao refrao_final se não existir
                        existentes = {_norm([l]) for l in refrao_final} if refrao_final else set()
                        for chave in chaves_refrao:
                            # Reconstrói as linhas a partir de um exemplo
                            exemplo = next((blk for blk in versos_processados if _norm(blk) == chave), None)
                            if exemplo:
                                texto_refrao = '\n'.join(exemplo)
                                if texto_refrao.strip() and _norm([texto_refrao]) not in existentes:
                                    refrao_final.append(texto_refrao)
                                    existentes.add(_norm([texto_refrao]))
                        # Remove as estrofes repetidas do conjunto de versos
                        versos_processados = [blk for blk in versos_processados if _norm(blk) not in chaves_refrao]
                        # Atualiza widgets imediatamente para refletir a detecção
                        try:
                            data_gui["refrao_widget"].delete("1.0", tk.END)
                            if refrao_final:
                                data_gui["refrao_widget"].insert(tk.END, '\n'.join(refrao_final))
                            data_gui["verso_widget"].delete("1.0", tk.END)
                            # Reescreve as estrofes restantes
                            if versos_processados:
                                texto_restante = '\n\n'.join(['\n'.join(blk) for blk in versos_processados])
                                data_gui["verso_widget"].insert(tk.END, texto_restante)
                        except Exception:
                            pass
            except Exception:
                pass
        elif nome_parte_gui.endswith("_Fixo"): 
            nome_original_para_default = nome_parte_gui.replace("_Fixo", "")
            defaults = self.DEFAULT_TEXTS.get(nome_original_para_default, {}); refrao_final = defaults.get("refrao", []); versos_processados = defaults.get("versos", [])
            tamanho_refrao, nome_refrao, bold_refrao, italic_refrao = DEFAULT_TAMANHO_FONTE_MUSICA_REFRAO, NOME_FONTE_PADRAO, True, False
            tamanho_verso, nome_verso, bold_verso, italic_verso = DEFAULT_TAMANHO_FONTE_MUSICA_VERSO, NOME_FONTE_PADRAO, True, False
            iniciar_com_refrao = False; aplicar_uppercase = True 
        if not refrao_final and not versos_processados: return False 
        if aplicar_uppercase:
            refrao_final = [s.upper() for s in refrao_final]; versos_processados = [[line.upper() for line in estrofe] for estrofe in versos_processados]
        titulo_adicionado = adiciona_texto_com_divisao(prs, layout_slide_branco, [titulo_parte], COR_TITULO, TAMANHO_TITULO_PARTE, NOME_FONTE_PADRAO, True, False, 5, use_auto_size=False)
        if titulo_adicionado: conteudo_adicionado_total = True
        if iniciar_com_refrao and refrao_final:
            if adiciona_texto_com_paginacao_inteligente(prs, layout_slide_branco, refrao_final, COR_REFRAO, tamanho_refrao, nome_refrao, bold_refrao, italic_refrao): conteudo_adicionado_total = True 
        for estrofe in versos_processados:
            if adiciona_texto_com_paginacao_inteligente(prs, layout_slide_branco, estrofe, COR_VERSO, tamanho_verso, nome_verso, bold_verso, italic_verso): conteudo_adicionado_total = True 
            if refrao_final: 
                if adiciona_texto_com_paginacao_inteligente(prs, layout_slide_branco, refrao_final, COR_REFRAO, tamanho_refrao, nome_refrao, bold_refrao, italic_refrao): conteudo_adicionado_total = True 
        if not versos_processados and refrao_final and not iniciar_com_refrao:
            if adiciona_texto_com_paginacao_inteligente(prs, layout_slide_branco, refrao_final, COR_REFRAO, tamanho_refrao, nome_refrao, bold_refrao, italic_refrao): conteudo_adicionado_total = True 
        return conteudo_adicionado_total

    def adicionar_leitura_slide_unico(self, prs, layout_slide_branco, nome_parte_gui):
        conteudo_adicionado = False; data_gui = self.widgets_gui.get(nome_parte_gui)
        if not data_gui or data_gui.get("tipo") != "leitura": return False
        titulo_amarelo_gui_str = data_gui["titulo_amarelo_widget"].get("1.0", tk.END).strip(); texto_branco_gui_str = data_gui["texto_branco_widget"].get("1.0", tk.END).strip()
        tamanho_titulo, nome_titulo, bold_titulo, italic_titulo = self._get_font_style_from_gui(data_gui, "titulo_amarelo", DEFAULT_TAMANHO_FONTE_LEITURA_TITULO_AMARELO, True, False)
        tamanho_texto, nome_texto, bold_texto, italic_texto = self._get_font_style_from_gui(data_gui, "texto_branco", DEFAULT_TAMANHO_FONTE_LEITURA_TEXTO_BRANCO, True, False)
        defaults = self.DEFAULT_TEXTS.get(nome_parte_gui, {}); titulo_amarelo_padrao = defaults.get("titulo_amarelo", []); texto_branco_padrao = defaults.get("texto_branco", [])
        titulo_amarelo_final = [l.strip() for l in titulo_amarelo_gui_str.split('\n') if l.strip()] if titulo_amarelo_gui_str else titulo_amarelo_padrao
        texto_branco_final = [l.strip() for l in texto_branco_gui_str.split('\n') if l.strip()] if texto_branco_gui_str else texto_branco_padrao
        if not titulo_amarelo_final and not texto_branco_final: return False
        slide = prs.slides.add_slide(layout_slide_branco); conteudo_adicionado = True
        esquerda=MARGEM_TEXTO; topo=MARGEM_TEXTO; largura=LARGURA_SLIDE-(2*MARGEM_TEXTO); altura=ALTURA_SLIDE-(2*MARGEM_TEXTO)
        caixa_texto = slide.shapes.add_textbox(esquerda,topo,largura,altura); frame_texto=caixa_texto.text_frame; frame_texto.clear(); frame_texto.word_wrap=True; frame_texto.vertical_anchor=MSO_ANCHOR.MIDDLE; frame_texto.auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        if titulo_amarelo_final:
            p_titulo = frame_texto.add_paragraph(); p_titulo.text = " ".join(titulo_amarelo_final); p_titulo.alignment = PP_ALIGN.CENTER
            p_titulo.font.name = nome_titulo; p_titulo.font.size = tamanho_titulo; p_titulo.font.color.rgb = COR_REFRAO; p_titulo.font.bold = bold_titulo; p_titulo.font.italic = italic_titulo
        if texto_branco_final:
            if titulo_amarelo_final: frame_texto.add_paragraph().text = ""
            # usar paginador inteligente para o texto principal
            # Fecha este slide e renderiza as páginas do texto
            try:
                caixa_texto.left=esquerda; caixa_texto.top=topo; caixa_texto.width=largura; caixa_texto.height=altura
            except Exception:
                pass
            # Cria páginas do texto em branco
            adiciona_texto_com_paginacao_inteligente(prs, layout_slide_branco, texto_branco_final, COR_VERSO, tamanho_texto, nome_texto, bold_texto, italic_texto)
            return True
        try: caixa_texto.left=esquerda; caixa_texto.top=topo; caixa_texto.width=largura; caixa_texto.height=altura; frame_texto.margin_bottom=Inches(0.05); frame_texto.margin_top=Inches(0.05); frame_texto.margin_left=Inches(0.1); frame_texto.margin_right=Inches(0.1)
        except Exception: pass
        return conteudo_adicionado

    def adicionar_aclamacao_slide_unico(self, prs, layout_slide_branco, nome_parte_gui):
        conteudo_adicionado = False; data_gui = self.widgets_gui.get(nome_parte_gui)
        if not data_gui or data_gui.get("tipo") != "aclamacao": return False
        titulo_secao = self.DEFAULT_TEXTS.get(nome_parte_gui, {}).get("titulo", nome_parte_gui.upper())
        if "titulo_secao_entry" in data_gui:
            gui_titulo = data_gui["titulo_secao_entry"].get().strip();
            if gui_titulo: titulo_secao = gui_titulo
        aclamacao_gui_str = data_gui["aclamacao_widget"].get("1.0", tk.END).strip(); antifona_gui_str = data_gui["antifona_widget"].get("1.0", tk.END).strip()
        tamanho_ac, nome_ac, bold_ac, italic_ac = self._get_font_style_from_gui(data_gui, "aclamacao", DEFAULT_TAMANHO_FONTE_ACLAMACAO, True, False)
        tamanho_an, nome_an, bold_an, italic_an = self._get_font_style_from_gui(data_gui, "antifona", DEFAULT_TAMANHO_FONTE_ANTIFONA, True, False)
        defaults = self.DEFAULT_TEXTS.get(nome_parte_gui, {}); aclamacao_padrao = defaults.get("aclamacao_texto", []); antifona_padrao = defaults.get("antifona_texto", [])
        aclamacao_final = [l.strip() for l in aclamacao_gui_str.split('\n') if l.strip()] if aclamacao_gui_str else aclamacao_padrao
        antifona_final = [l.strip() for l in antifona_gui_str.split('\n') if l.strip()] if antifona_gui_str else antifona_padrao
        if not aclamacao_final and not antifona_final: return False
        aplicar_uppercase = True
        if "uppercase_var" in data_gui:
            try: aplicar_uppercase = data_gui["uppercase_var"].get()
            except Exception: aplicar_uppercase = True
        if aplicar_uppercase: aclamacao_final = [s.upper() for s in aclamacao_final]; antifona_final = [s.upper() for s in antifona_final]
        if adiciona_texto_com_divisao(prs, layout_slide_branco, [titulo_secao], COR_TITULO, TAMANHO_TITULO_PARTE, NOME_FONTE_PADRAO, True, False, 5, use_auto_size=False): conteudo_adicionado = True
        slide = prs.slides.add_slide(layout_slide_branco); conteudo_adicionado = True 
        esquerda=MARGEM_TEXTO; topo=MARGEM_TEXTO; largura=LARGURA_SLIDE-(2*MARGEM_TEXTO); altura=ALTURA_SLIDE-(2*MARGEM_TEXTO)
        caixa_texto = slide.shapes.add_textbox(esquerda,topo,largura,altura); frame_texto=caixa_texto.text_frame; frame_texto.clear(); frame_texto.word_wrap=True; frame_texto.vertical_anchor=MSO_ANCHOR.MIDDLE; frame_texto.auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        if aclamacao_final:
            p_ac = frame_texto.add_paragraph(); p_ac.text = " ".join(aclamacao_final); p_ac.alignment = PP_ALIGN.CENTER
            p_ac.font.name=nome_ac; p_ac.font.size=tamanho_ac; p_ac.font.color.rgb=COR_REFRAO; p_ac.font.bold=bold_ac; p_ac.font.italic=italic_ac
        if antifona_final:
            # Paginar antifona com inteligência
            adiciona_texto_com_paginacao_inteligente(prs, layout_slide_branco, antifona_final, COR_VERSO, tamanho_an, nome_an, bold_an, italic_an)
            return True
        try: caixa_texto.left=esquerda; caixa_texto.top=topo; caixa_texto.width=largura; caixa_texto.height=altura; frame_texto.margin_bottom=Inches(0.05); frame_texto.margin_top=Inches(0.05); frame_texto.margin_left=Inches(0.1); frame_texto.margin_right=Inches(0.1)
        except Exception: pass
        return conteudo_adicionado 

    def adicionar_secao_fixa(self, prs, layout_slide_branco, titulo_secao, texto_linhas, tamanho_fonte, linhas_por_slide_custom, cor=COR_VERSO, bold_content=True, use_auto_size_content=False):
        titulo_adicionado = False; conteudo_adicionado_slides = False
        if titulo_secao and titulo_secao.strip():
            if adiciona_texto_com_divisao(prs, layout_slide_branco, [titulo_secao], COR_TITULO, TAMANHO_TITULO_PARTE, NOME_FONTE_PADRAO, True, False, 5, use_auto_size=False): titulo_adicionado = True
        if texto_linhas:
            if adiciona_texto_com_divisao(prs, layout_slide_branco, texto_linhas, cor, tamanho_fonte, NOME_FONTE_PADRAO, bold_content, False, linhas_por_slide_custom, use_auto_size=use_auto_size_content): conteudo_adicionado_slides = True
        return titulo_adicionado or conteudo_adicionado_slides

    def adicionar_secao_palavra(self, prs, layout_slide_branco, nome_parte_gui):
        conteudo_adicionado_total = False; data_gui = self.widgets_gui.get(nome_parte_gui)
        if not data_gui or data_gui.get("tipo") != "palavra": return False
        titulo_secao = self.DEFAULT_TEXTS.get(nome_parte_gui, {}).get("titulo", nome_parte_gui.upper())
        if "titulo_secao_entry" in data_gui:
            gui_titulo = data_gui["titulo_secao_entry"].get().strip();
            if gui_titulo: titulo_secao = gui_titulo
        texto_gui_str = data_gui["texto_widget"].get("1.0", tk.END).strip()
        tamanho_fonte, nome_fonte, bold_state, italic_state = self._get_font_style_from_gui(data_gui, "texto", DEFAULT_TAMANHO_FONTE_PALAVRA, True, False)
        texto_padrao = self.DEFAULT_TEXTS.get(nome_parte_gui, {}).get("texto", [])
        texto_final = [l.strip() for l in texto_gui_str.split('\n') if l.strip()] if texto_gui_str else texto_padrao
        if not texto_final: return False
        aplicar_uppercase = True
        if "uppercase_var" in data_gui:
            try: aplicar_uppercase = data_gui["uppercase_var"].get()
            except Exception: aplicar_uppercase = True
        if aplicar_uppercase: texto_final = [s.upper() for s in texto_final]
        if adiciona_texto_com_divisao(prs, layout_slide_branco, [titulo_secao], COR_TITULO, TAMANHO_TITULO_PARTE, NOME_FONTE_PADRAO, True, False, 5, use_auto_size=False): conteudo_adicionado_total = True
        if adiciona_texto_com_divisao(prs, layout_slide_branco, texto_final, COR_TITULO, tamanho_fonte, nome_fonte, bold_state, italic_state, LINHAS_POR_SLIDE_PALAVRA, use_auto_size=True): conteudo_adicionado_total = True 
        return conteudo_adicionado_total

    def adicionar_aviso_com_imagem(self, prs, layout_slide_branco, nome_arquivo_imagem):
        slide_adicionado = False
        if getattr(sys, 'frozen', False) and hasattr(sys, '_MEIPASS'): application_path = sys._MEIPASS
        else:
            try: application_path = os.path.dirname(os.path.abspath(__file__))
            except NameError: application_path = os.getcwd()
        caminho_imagem = os.path.join(application_path, nome_arquivo_imagem)
        if os.path.exists(caminho_imagem):
            try:
                slide_avisos = prs.slides.add_slide(layout_slide_branco); slide_adicionado = True
                slide_avisos.shapes.add_picture(caminho_imagem, Inches(0), Inches(0), width=LARGURA_SLIDE, height=ALTURA_SLIDE)
            except Exception as e_img:
                messagebox.showerror("Erro Imagem Avisos", f"Não foi possível adicionar a imagem de avisos:\n{e_img}", parent=self.master)
                if adiciona_texto_com_divisao(prs, layout_slide_branco, ["AVISOS"], COR_TITULO, TAMANHO_TITULO_PARTE, NOME_FONTE_PADRAO, True, False, 5, use_auto_size=False): slide_adicionado = True
        else:
            messagebox.showwarning("Imagem Avisos Não Encontrada", f"O arquivo '{nome_arquivo_imagem}' não foi encontrado em '{application_path}'.\nVerifique se ele está na mesma pasta do executável/script.", parent=self.master)
            if adiciona_texto_com_divisao(prs, layout_slide_branco, ["AVISOS"], COR_TITULO, TAMANHO_TITULO_PARTE, NOME_FONTE_PADRAO, True, False, 5, use_auto_size=False): slide_adicionado = True
        return slide_adicionado

    def gerar_apresentacao(self):
        self.status_label.config(text="Gerando apresentação...")
        self.master.update_idletasks()
        self._reconstruir_ordem_geracao_dinamica() 

        try:
            prs = Presentation()
            prs.slide_width = LARGURA_SLIDE; prs.slide_height = ALTURA_SLIDE
            try:
                fill = prs.slide_masters[0].background.fill
                fill.solid(); fill.fore_color.rgb = COR_FUNDO_PRETO
            except Exception as e_bg_master: print(f"Aviso: Não foi possível aplicar fundo preto no master: {e_bg_master}")

            layout_slide_branco = next((l for l in prs.slide_layouts if "Branco" in l.name or "Blank" in l.name), None)
            if not layout_slide_branco: 
                idx_fallback = 5 if len(prs.slide_layouts) > 5 else (len(prs.slide_layouts) -1 if len(prs.slide_layouts) > 0 else 0)
                if len(prs.slide_layouts) > 0: layout_slide_branco = prs.slide_layouts[idx_fallback]
                else: 
                    sl = prs.slide_masters[0].slide_layouts.add_layout("Branco Personalizado", prs.slide_masters[0])
                    layout_slide_branco = sl
            
            slides_adicionados_conteudo_geral = 0 
            
            for i, nome_parte in enumerate(self.ordem_geracao_dinamica):
                conteudo_desta_secao_adicionado = False 
                
                if nome_parte == "TITULO_INICIAL_PLACEHOLDER":
                    initial_title_str = self.initial_title_widget.get("1.0", tk.END).strip()
                    initial_title_lines = [l.strip() for l in initial_title_str.split('\n') if l.strip()]
                    if initial_title_lines:
                        if adiciona_texto_com_divisao(prs, layout_slide_branco, initial_title_lines, COR_TITULO, TAMANHO_FONTE_TITULO_INICIAL, NOME_FONTE_PADRAO, True, False, 5, use_auto_size=True):
                            conteudo_desta_secao_adicionado = True
                elif nome_parte == "CREDO":
                    if self.adicionar_secao_fixa(prs, layout_slide_branco, "ORAÇÃO DO CREDO", TEXTO_CREDO, Pt(90), 3, use_auto_size_content=True):
                        conteudo_desta_secao_adicionado = True
                elif nome_parte == "PRECES":
                    if self.adicionar_secao_fixa(prs, layout_slide_branco, "PRECES", [], Pt(1), 1): 
                        conteudo_desta_secao_adicionado = True
                elif nome_parte == "SANTO_TITULO":
                    if self.adicionar_secao_fixa(prs, layout_slide_branco, "SANTO", [], Pt(1),1):
                        conteudo_desta_secao_adicionado = True
                elif nome_parte == "ORACAO_EUCARISTICA":
                    if self.adicionar_secao_fixa(prs, layout_slide_branco, "ORAÇÃO EUCARÍSTICA", [], Pt(1), 2):
                        conteudo_desta_secao_adicionado = True
                    # Se "Selecionar todas" estiver ativo, inserir um slide-título com o nome de cada oração antes do seu conteúdo
                    if self.eucaristica_selecionar_todas.get():
                        for nome, bloco in self._obter_oracoes_eucaristicas_com_nomes():
                            titulo_ok = self.adicionar_secao_fixa(prs, layout_slide_branco, nome.upper(), [], Pt(1), 1)
                            conteudo_desta_secao_adicionado = conteudo_desta_secao_adicionado or titulo_ok
                            linhas = [(linha or "").upper() for linha in bloco]
                            if linhas:
                                if self.adicionar_secao_fixa(prs, layout_slide_branco, "", linhas, Pt(self.eucaristica_font_size_pt), 1, cor=COR_REFRAO, use_auto_size_content=True):
                                    conteudo_desta_secao_adicionado = True
                    else:
                        conteudo_oracoes = self._obter_oracoes_eucaristicas_para_geracao()
                        if conteudo_oracoes:
                            linhas = []
                            for bloco in conteudo_oracoes:
                                for linha in bloco:
                                    linhas.append((linha or "").upper())
                                # separador entre orações diferentes (não cria slide vazio, só força quebra posterior)
                                linhas.append("")
                            while linhas and linhas[-1] == "":
                                linhas.pop()
                            # Texto em amarelo e 1 linha por slide
                            if self.adicionar_secao_fixa(prs, layout_slide_branco, "", linhas, Pt(self.eucaristica_font_size_pt), 1, cor=COR_REFRAO, use_auto_size_content=True):
                                conteudo_desta_secao_adicionado = True
    
                elif nome_parte == "CORDEIRO_TITULO":
                    if self.adicionar_secao_fixa(prs, layout_slide_branco, "CORDEIRO", [], Pt(1), 1):
                        conteudo_desta_secao_adicionado = True
                elif nome_parte == "SANTA_LUZIA":
                    if self.adicionar_secao_fixa(prs, layout_slide_branco, "ORAÇÃO A SANTA LUZIA", TEXTO_ORACAO_SANTA_LUZIA, Pt(90), 4, use_auto_size_content=True):
                         conteudo_desta_secao_adicionado = True
                elif nome_parte == "AVISOS_IMG": 
                    if self.adicionar_aviso_com_imagem(prs, layout_slide_branco, "AVISOS.png"):
                         conteudo_desta_secao_adicionado = True
                elif nome_parte in self.widgets_gui: 
                    data_gui = self.widgets_gui[nome_parte]
                    tipo_secao = data_gui.get("tipo")
                    if tipo_secao == "musica":
                        if self.adicionar_secao_musical(prs, layout_slide_branco, nome_parte): conteudo_desta_secao_adicionado = True
                    elif tipo_secao == "leitura":
                        if self.adicionar_leitura_slide_unico(prs, layout_slide_branco, nome_parte): conteudo_desta_secao_adicionado = True
                    elif tipo_secao == "aclamacao":
                        if self.adicionar_aclamacao_slide_unico(prs, layout_slide_branco, nome_parte): conteudo_desta_secao_adicionado = True
                    elif tipo_secao == "palavra":
                        if self.adicionar_secao_palavra(prs, layout_slide_branco, nome_parte): conteudo_desta_secao_adicionado = True
                
                if conteudo_desta_secao_adicionado:
                    slides_adicionados_conteudo_geral += 1
                    is_last_item_real_na_apresentacao = (nome_parte == "AVISOS_IMG") 
                    
                    if not is_last_item_real_na_apresentacao:
                        if len(prs.slides) > 0: 
                            last_slide_is_separator = not prs.slides[-1].shapes and not prs.slides[-1].placeholders
                            if not last_slide_is_separator:
                                prs.slides.add_slide(layout_slide_branco)
                        elif slides_adicionados_conteudo_geral > 0 : 
                             prs.slides.add_slide(layout_slide_branco)
            
            if slides_adicionados_conteudo_geral > 0 and len(prs.slides) > 0:
                if not (len(prs.slides) == 1 and self.initial_title_widget.get("1.0", tk.END).strip() and slides_adicionados_conteudo_geral == 1) :
                    last_slide = prs.slides[-1]
                    if not last_slide.shapes and not last_slide.placeholders and \
                       (not self.ordem_geracao_dinamica or self.ordem_geracao_dinamica[-1] != "AVISOS_IMG" or nome_parte != "AVISOS_IMG"): 
                         xml_slides = prs.slides._sldIdLst; slides = list(xml_slides)
                         if slides: xml_slides.remove(slides[-1])

            if slides_adicionados_conteudo_geral == 0 and not (self.initial_title_widget.get("1.0", tk.END).strip()):
                if len(prs.slides) > 0: prs.slides._sldIdLst.clear()

            if not prs.slides: 
                messagebox.showwarning("Atenção", "Nenhum conteúdo resultou em slides. O arquivo não será salvo.", parent=self.master)
                self.status_label.config(text="Geração cancelada (vazia).")
                return

            filepath = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint Presentations", "*.pptx")], initialfile="Missa_Gerada_v35_2.pptx", parent=self.master)
            if not filepath: self.status_label.config(text="Geração cancelada."); return
            prs.save(filepath)
            self.status_label.config(text=f"Salvo: {os.path.basename(filepath)}")
            messagebox.showinfo("Sucesso", f"Apresentação '{os.path.basename(filepath)}' gerada e salva!", parent=self.master)
            try:
                if platform.system() == 'Darwin': subprocess.call(('open', filepath))
                elif platform.system() == 'Windows': os.startfile(filepath)
                else: subprocess.call(('xdg-open', filepath))
            except Exception as e_open: print(f"Erro ao abrir o arquivo: {e_open}")

        except Exception as e:
            self.status_label.config(text="Erro durante a geração!")
            import traceback; traceback.print_exc()
            messagebox.showerror("Erro", f"Ocorreu um erro:\n{e}", parent=self.master)
        finally:
            self.master.update_idletasks()

if __name__ == "__main__":
    root = tk.Tk() # Revertido para tk.Tk()
    app = MassSlideGeneratorApp(root)
    root.mainloop()